from django.shortcuts import render
from django.http import JsonResponse, HttpResponse
from django.db.models import Min, Max, Avg
from django.views.decorators.csrf import csrf_exempt
from decimal import Decimal
import json
import csv
import io
from datetime import datetime, timedelta
from .data import (
    FCP_FICHE_SIGNALETIQUE, 
    get_all_fcp_names, 
    get_fcp_data, 
    get_risk_label,
    get_type_icon,
    get_type_color
)
from .models import FicheSignaletique, FCP_VL_MODELS, get_vl_model

# ============================================================================
# CONSTANTES FINANCIÈRES
# ============================================================================
# Nombre de jours de trading par an (standard pour les marchés financiers)
# Le marché BRVM est ouvert ~252 jours par an (hors week-ends et jours fériés)
TRADING_DAYS_PER_YEAR = 365

# Taux sans risque annuel (en %) - Taux de référence BCEAO
RISK_FREE_RATE_ANNUAL = 3.25

# Taux sans risque journalier
RISK_FREE_RATE_DAILY = RISK_FREE_RATE_ANNUAL / TRADING_DAYS_PER_YEAR

# Facteur d'annualisation pour la volatilité
ANNUALIZATION_FACTOR = TRADING_DAYS_PER_YEAR ** 0.5
# ============================================================================

# Create your views here.


def valeurs_liquidatives(request):
    """Vue pour la page des Valeurs Liquidatives"""
    # Récupérer le FCP sélectionné
    selected_fcp = request.GET.get('fcp', 'FCP PLACEMENT AVANTAGE')
    
    # Liste des FCP depuis la base
    fcp_list = list(FicheSignaletique.objects.values_list('nom', flat=True).order_by('nom'))
    
    # Récupérer les données VL pour le FCP sélectionné
    vl_model = get_vl_model(selected_fcp)
    vl_data = []
    stats = {}
    perf_calendaires = {}
    perf_glissantes = {}
    analyse_stats = {}
    tracking_error = {}
    
    if vl_model:
        # Récupérer toutes les VL ordonnées par date
        vl_queryset = vl_model.objects.all().order_by('date')
        
        # Convertir en liste pour JSON
        for vl in vl_queryset:
            vl_data.append({
                'date': vl.date.strftime('%Y-%m-%d'),
                'valeur': float(vl.valeur)
            })
        
        # Calculer les statistiques
        if vl_queryset.exists():
            latest_vl = vl_queryset.last()
            first_vl = vl_queryset.first()
            today = latest_vl.date  # Dernière date de la base (pas la date du jour)
            
            # Helper pour calculer la performance
            def calc_perf(vl_ref):
                if vl_ref:
                    return round(((float(latest_vl.valeur) / float(vl_ref.valeur)) - 1) * 100, 2)
                return None
            
            # VL pour différentes périodes GLISSANTES (depuis la dernière date en base)
            vl_1d = vl_queryset.filter(date__lt=today).last()
            vl_1m = vl_queryset.filter(date__lte=today - timedelta(days=30)).last()
            vl_3m = vl_queryset.filter(date__lte=today - timedelta(days=90)).last()
            vl_6m = vl_queryset.filter(date__lte=today - timedelta(days=180)).last()
            vl_1y = vl_queryset.filter(date__lte=today - timedelta(days=365)).last()
            vl_3y = vl_queryset.filter(date__lte=today - timedelta(days=365*3)).last()
            vl_5y = vl_queryset.filter(date__lte=today - timedelta(days=365*5)).last()
            
            # Performances calendaires (To-Date) - basées sur la dernière date en base
            # WTD: dernière VL de la semaine précédente (vendredi ou avant le lundi de cette semaine)
            start_of_week = today - timedelta(days=today.weekday())  # Lundi de la semaine
            vl_wtd = vl_queryset.filter(date__lt=start_of_week).last()
            
            # MTD: dernière VL du mois précédent (dernier jour avant le 1er du mois)
            start_of_month = today.replace(day=1)
            vl_mtd = vl_queryset.filter(date__lt=start_of_month).last()
            
            # QTD: dernière VL du trimestre précédent
            quarter_month = ((today.month - 1) // 3) * 3 + 1
            start_of_quarter = today.replace(month=quarter_month, day=1)
            vl_qtd = vl_queryset.filter(date__lt=start_of_quarter).last()
            
            # STD: dernière VL du semestre précédent
            semester_month = 1 if today.month <= 6 else 7
            start_of_semester = today.replace(month=semester_month, day=1)
            vl_std = vl_queryset.filter(date__lt=start_of_semester).last()
            
            # YTD: dernière VL de l'année précédente (dernier jour avant le 1er janvier)
            start_of_year = today.replace(month=1, day=1)
            vl_ytd = vl_queryset.filter(date__lt=start_of_year).last()
            
            perf_calendaires = {
                'wtd': calc_perf(vl_wtd),
                'mtd': calc_perf(vl_mtd),
                'qtd': calc_perf(vl_qtd),
                'std': calc_perf(vl_std),
                'ytd': calc_perf(vl_ytd),
            }
            
            # Performances glissantes
            perf_glissantes = {
                'perf_1m': calc_perf(vl_1m),
                'perf_3m': calc_perf(vl_3m),
                'perf_6m': calc_perf(vl_6m),
                'perf_1y': calc_perf(vl_1y),
                'perf_3y': calc_perf(vl_3y),
                'perf_5y': calc_perf(vl_5y),
                'origine': calc_perf(first_vl),
            }
            
            stats = {
                'derniere_vl': float(latest_vl.valeur),
                'derniere_date': latest_vl.date.strftime('%d/%m/%Y'),
                'premiere_vl': float(first_vl.valeur),
                'premiere_date': first_vl.date.strftime('%d/%m/%Y'),
                'var_1j': calc_perf(vl_1d) or 0,
                'var_1m': calc_perf(vl_1m) or 0,
                'var_1y': calc_perf(vl_1y) or 0,
                'var_ytd': perf_calendaires['ytd'] or 0,
                'var_origine': perf_glissantes['origine'] or 0,
                'nb_vl': vl_queryset.count(),
            }
            
            # Calculer la Tracking Error (volatilité annualisée) pour différentes périodes
            def calc_tracking_error(start_date):
                """Calcule la tracking error (volatilité annualisée) depuis une date de référence"""
                if start_date is None:
                    return None
                period_vl = list(vl_queryset.filter(date__gte=start_date).values_list('valeur', flat=True))
                if len(period_vl) < 2:
                    return None
                period_valeurs = [float(v) for v in period_vl]
                period_rendements = [(period_valeurs[i] / period_valeurs[i-1] - 1) * 100 for i in range(1, len(period_valeurs))]
                if len(period_rendements) < 2:
                    return None
                moyenne = sum(period_rendements) / len(period_rendements)
                variance = sum((r - moyenne) ** 2 for r in period_rendements) / (len(period_rendements) - 1)  # Variance non biaisée
                ecart_type = variance ** 0.5
                # Annualiser la volatilité
                volatilite_ann = ecart_type * ANNUALIZATION_FACTOR
                return round(volatilite_ann, 2)
            
            tracking_error = {
                'wtd': calc_tracking_error(start_of_week),
                'mtd': calc_tracking_error(start_of_month),
                'qtd': calc_tracking_error(start_of_quarter),
                'std': calc_tracking_error(start_of_semester),
                'ytd': calc_tracking_error(start_of_year),
                'origine': calc_tracking_error(first_vl.date if first_vl else None),
            }
            
            # Calculer les statistiques pour l'onglet Analyse
            valeurs = [float(v) for v in vl_queryset.values_list('valeur', flat=True)]
            if len(valeurs) > 1:
                rendements = [(valeurs[i] / valeurs[i-1] - 1) * 100 for i in range(1, len(valeurs))]
                
                # Statistiques descriptives
                moyenne_rdt = sum(rendements) / len(rendements)
                variance = sum((r - moyenne_rdt) ** 2 for r in rendements) / (len(rendements) - 1)  # Variance non biaisée
                ecart_type = variance ** 0.5
                volatilite_ann = ecart_type * ANNUALIZATION_FACTOR
                
                rendements_sorted = sorted(rendements)
                n = len(rendements_sorted)
                mediane = rendements_sorted[n // 2] if n % 2 == 1 else (rendements_sorted[n//2 - 1] + rendements_sorted[n//2]) / 2
                
                # Min/Max
                vl_min = min(valeurs)
                vl_max = max(valeurs)
                rdt_min = min(rendements)
                rdt_max = max(rendements)
                
                # Jours positifs/négatifs
                jours_positifs = sum(1 for r in rendements if r > 0)
                jours_negatifs = sum(1 for r in rendements if r < 0)
                
                # Drawdown max
                peak = valeurs[0]
                max_drawdown = 0
                for v in valeurs:
                    if v > peak:
                        peak = v
                    drawdown = (peak - v) / peak * 100
                    if drawdown > max_drawdown:
                        max_drawdown = drawdown
                
                # Ratio de Sharpe (avec taux sans risque constant)
                sharpe = ((moyenne_rdt - RISK_FREE_RATE_DAILY) / ecart_type * ANNUALIZATION_FACTOR) if ecart_type > 0 else 0
                
                # Ratio de Sortino
                rendements_negatifs = [r for r in rendements if r < 0]
                if rendements_negatifs:
                    downside_var = sum(r ** 2 for r in rendements_negatifs) / (len(rendements_negatifs) - 1) if len(rendements_negatifs) > 1 else sum(r ** 2 for r in rendements_negatifs)
                    downside_std = downside_var ** 0.5
                    sortino = ((moyenne_rdt - RISK_FREE_RATE_DAILY) / downside_std * ANNUALIZATION_FACTOR) if downside_std > 0 else 0
                else:
                    sortino = 0
                
                # Skewness (asymétrie)
                if ecart_type > 0:
                    skewness = sum((r - moyenne_rdt) ** 3 for r in rendements) / (n * ecart_type ** 3)
                else:
                    skewness = 0
                
                # Kurtosis (aplatissement) - excess kurtosis (0 pour normal)
                if ecart_type > 0:
                    kurtosis = sum((r - moyenne_rdt) ** 4 for r in rendements) / (n * ecart_type ** 4) - 3
                else:
                    kurtosis = 0
                
                # VaR Historique (Value at Risk)
                var_95 = rendements_sorted[int(n * 0.05)]  # 5ème percentile
                var_99 = rendements_sorted[int(n * 0.01)]  # 1er percentile
                
                # CVaR / Expected Shortfall (moyenne des pertes au-delà de la VaR)
                cvar_95_values = [r for r in rendements if r <= var_95]
                cvar_99_values = [r for r in rendements if r <= var_99]
                cvar_95 = sum(cvar_95_values) / len(cvar_95_values) if cvar_95_values else var_95
                cvar_99 = sum(cvar_99_values) / len(cvar_99_values) if cvar_99_values else var_99
                
                # Calcul des drawdowns détaillés
                drawdowns_data = []
                dates_list = list(vl_queryset.values_list('date', flat=True))
                peak = valeurs[0]
                peak_date = dates_list[0]
                current_dd_start = None
                current_dd_peak = valeurs[0]
                
                underwater_data = []  # Pour le graphique underwater
                
                for i, v in enumerate(valeurs):
                    if v > peak:
                        # Nouveau pic atteint
                        if current_dd_start is not None:
                            # Fin du drawdown précédent
                            dd_depth = (current_dd_peak - min(valeurs[current_dd_start:i])) / current_dd_peak * 100
                            dd_trough_idx = current_dd_start + valeurs[current_dd_start:i].index(min(valeurs[current_dd_start:i]))
                            drawdowns_data.append({
                                'start_date': dates_list[current_dd_start].strftime('%Y-%m-%d'),
                                'trough_date': dates_list[dd_trough_idx].strftime('%Y-%m-%d'),
                                'end_date': dates_list[i].strftime('%Y-%m-%d'),
                                'depth': round(dd_depth, 2),
                                'duration': i - current_dd_start,
                                'recovery': i - dd_trough_idx
                            })
                            current_dd_start = None
                        peak = v
                        peak_date = dates_list[i]
                        current_dd_peak = v
                    else:
                        if current_dd_start is None and v < peak:
                            current_dd_start = i - 1 if i > 0 else 0
                            current_dd_peak = peak
                    
                    # Calculer le drawdown courant pour underwater chart
                    dd_current = (peak - v) / peak * 100
                    underwater_data.append({
                        'date': dates_list[i].strftime('%Y-%m-%d'),
                        'drawdown': round(-dd_current, 2)  # Négatif pour l'affichage
                    })
                
                # Si on est encore en drawdown à la fin
                if current_dd_start is not None:
                    dd_depth = (current_dd_peak - min(valeurs[current_dd_start:])) / current_dd_peak * 100
                    dd_trough_idx = current_dd_start + valeurs[current_dd_start:].index(min(valeurs[current_dd_start:]))
                    drawdowns_data.append({
                        'start_date': dates_list[current_dd_start].strftime('%Y-%m-%d'),
                        'trough_date': dates_list[dd_trough_idx].strftime('%Y-%m-%d'),
                        'end_date': None,  # Encore en cours
                        'depth': round(dd_depth, 2),
                        'duration': len(valeurs) - current_dd_start,
                        'recovery': None  # Pas encore récupéré
                    })
                
                # Trier les drawdowns par profondeur (les pires en premier)
                drawdowns_data.sort(key=lambda x: x['depth'], reverse=True)
                
                # Histogramme des rendements (bins)
                hist_min = min(rendements)
                hist_max = max(rendements)
                nb_bins = 30
                bin_width = (hist_max - hist_min) / nb_bins if hist_max != hist_min else 1
                histogram_data = []
                for i in range(nb_bins):
                    bin_start = hist_min + i * bin_width
                    bin_end = bin_start + bin_width
                    count = sum(1 for r in rendements if bin_start <= r < bin_end)
                    histogram_data.append({
                        'bin_start': round(bin_start, 3),
                        'bin_end': round(bin_end, 3),
                        'count': count,
                        'frequency': round(count / n * 100, 2)
                    })
                
                stats['volatilite'] = round(volatilite_ann, 2)
                
                analyse_stats = {
                    # Statistiques descriptives
                    'nb_observations': len(rendements),
                    'rendement_moyen': round(moyenne_rdt, 4),
                    'rendement_moyen_ann': round(moyenne_rdt * 365, 2),
                    'mediane': round(mediane, 4),
                    'ecart_type': round(ecart_type, 4),
                    'volatilite_ann': round(volatilite_ann, 2),
                    'vl_min': round(vl_min, 2),
                    'vl_max': round(vl_max, 2),
                    'rdt_min': round(rdt_min, 2),
                    'rdt_max': round(rdt_max, 2),
                    # Profil de risque
                    'max_drawdown': round(max_drawdown, 2),
                    'sharpe': round(sharpe, 2),
                    'sortino': round(sortino, 2),
                    'jours_positifs': jours_positifs,
                    'jours_negatifs': jours_negatifs,
                    'ratio_positif': round(jours_positifs / len(rendements) * 100, 1) if rendements else 0,
                    # Distribution des rendements
                    'skewness': round(skewness, 3),
                    'kurtosis': round(kurtosis, 3),
                    'var_95': round(var_95, 3),
                    'var_99': round(var_99, 3),
                    'cvar_95': round(cvar_95, 3),
                    'cvar_99': round(cvar_99, 3),
                    # Données pour graphiques (JSON)
                    'histogram_data': histogram_data,
                    'underwater_data': underwater_data,
                    'drawdowns_data': drawdowns_data[:10],  # Top 10 drawdowns
                    'rendements_list': [round(r, 4) for r in rendements],  # Pour histogramme JS
                }
            else:
                stats['volatilite'] = 0
                analyse_stats = {}
    
    # Préparer les données pour tous les FCP (pour le scatter plot)
    all_fcp_stats = []
    for fcp_name in fcp_list:
        fcp_vl_model = get_vl_model(fcp_name)
        if fcp_vl_model:
            fcp_vl = fcp_vl_model.objects.all().order_by('date')
            if fcp_vl.exists() and fcp_vl.count() > 30:
                first = fcp_vl.first()
                last = fcp_vl.last()
                rendement = ((float(last.valeur) / float(first.valeur)) - 1) * 100
                
                # Calculer volatilité
                valeurs = list(fcp_vl.values_list('valeur', flat=True))
                rendements = [(float(valeurs[i]) / float(valeurs[i-1]) - 1) for i in range(1, len(valeurs))]
                if rendements and len(rendements) > 1:
                    moyenne = sum(rendements) / len(rendements)
                    variance = sum((r - moyenne) ** 2 for r in rendements) / (len(rendements) - 1)  # Variance non biaisée
                    vol = (variance ** 0.5) * ANNUALIZATION_FACTOR * 100
                else:
                    vol = 0
                
                # Récupérer le type de fond depuis FicheSignaletique
                try:
                    fiche = FicheSignaletique.objects.get(nom=fcp_name)
                    type_fond = fiche.type_fond
                except FicheSignaletique.DoesNotExist:
                    type_fond = 'Diversifié'
                
                all_fcp_stats.append({
                    'nom': fcp_name,
                    'rendement': round(rendement, 2),
                    'volatilite': round(vol, 2),
                    'type_fond': type_fond,
                    'selected': fcp_name == selected_fcp
                })
    
    # Extraire les données pour les graphiques (JSON)
    histogram_data_json = json.dumps(analyse_stats.get('histogram_data', []))
    underwater_data_json = json.dumps(analyse_stats.get('underwater_data', []))
    drawdowns_data_json = json.dumps(analyse_stats.get('drawdowns_data', []))
    rendements_list_json = json.dumps(analyse_stats.get('rendements_list', []))
    
    # Données de fiche signalétique pour l'onglet Fiche signalétique
    selected_fcp_fiche_data = get_fcp_data(selected_fcp)
    fcp_enriched = {}
    for name, data in FCP_FICHE_SIGNALETIQUE.items():
        fcp_enriched[name] = {
            **data,
            'risk_label': get_risk_label(data['echelle_risque']),
            'type_icon': get_type_icon(data['type_fond']),
            'type_color': get_type_color(data['type_fond']),
        }
    
    selected_fcp_data = None
    if selected_fcp_fiche_data:
        selected_fcp_data = {
            **selected_fcp_fiche_data,
            'risk_label': get_risk_label(selected_fcp_fiche_data['echelle_risque']),
            'type_icon': get_type_icon(selected_fcp_fiche_data['type_fond']),
            'type_color': get_type_color(selected_fcp_fiche_data['type_fond']),
        }
    
    context = {
        'page_title': 'Valeurs Liquidatives',
        'page_description': 'Suivi et historique des valeurs liquidatives des FCP',
        'fcp_list': fcp_list,
        'selected_fcp': selected_fcp,
        'vl_data_json': json.dumps(vl_data),
        'stats': stats,
        'perf_calendaires': perf_calendaires,
        'perf_glissantes': perf_glissantes,
        'analyse_stats': analyse_stats,
        'tracking_error': tracking_error,
        'all_fcp_stats_json': json.dumps(all_fcp_stats),
        'histogram_data_json': histogram_data_json,
        'underwater_data_json': underwater_data_json,
        'drawdowns_data_json': drawdowns_data_json,
        'rendements_list_json': rendements_list_json,
        # Données fiche signalétique
        'fcp_data': fcp_enriched,
        'selected_fcp_data': selected_fcp_data,
        'fcp_data_json': json.dumps(fcp_enriched),
    }
    return render(request, 'fcp_app/valeurs_liquidatives.html', context)


def api_scatter_data(request):
    """API pour récupérer les données du scatter plot avec filtre de période"""
    period = request.GET.get('period', 'origin')
    
    fcp_list = list(FicheSignaletique.objects.values_list('nom', flat=True).order_by('nom'))
    all_fcp_stats = []
    
    for fcp_name in fcp_list:
        fcp_vl_model = get_vl_model(fcp_name)
        if fcp_vl_model:
            fcp_vl = fcp_vl_model.objects.all().order_by('date')
            if fcp_vl.exists() and fcp_vl.count() > 30:
                last = fcp_vl.last()
                latest_date = last.date  # Dernière date disponible en base
                
                # Pour les calculs To-Date, on cherche la dernière VL avant le début de la période
                if period == 'wtd':  # Week-to-Date (depuis le lundi de cette semaine)
                    start_of_week = latest_date - timedelta(days=latest_date.weekday())
                    ref_vl = fcp_vl.filter(date__lt=start_of_week).last()
                elif period == 'mtd':  # Month-to-Date (depuis le début du mois)
                    start_of_month = latest_date.replace(day=1)
                    ref_vl = fcp_vl.filter(date__lt=start_of_month).last()
                elif period == 'qtd':  # Quarter-to-Date (depuis le début du trimestre)
                    quarter_month = ((latest_date.month - 1) // 3) * 3 + 1
                    start_of_quarter = latest_date.replace(month=quarter_month, day=1)
                    ref_vl = fcp_vl.filter(date__lt=start_of_quarter).last()
                elif period == 'std':  # Semester-to-Date (depuis le début du semestre)
                    semester_month = 1 if latest_date.month <= 6 else 7
                    start_of_semester = latest_date.replace(month=semester_month, day=1)
                    ref_vl = fcp_vl.filter(date__lt=start_of_semester).last()
                elif period == 'ytd':  # Year-to-Date (depuis le début de l'année)
                    start_of_year = latest_date.replace(month=1, day=1)
                    ref_vl = fcp_vl.filter(date__lt=start_of_year).last()
                else:  # origin - depuis le début des données
                    ref_vl = fcp_vl.first()
                
                if ref_vl:
                    # Calculer le rendement
                    rendement = ((float(last.valeur) / float(ref_vl.valeur)) - 1) * 100
                    
                    # Calculer la volatilité sur la période sélectionnée
                    if period == 'origin':
                        filtered_vl = fcp_vl
                    else:
                        # Filtrer depuis la date de référence jusqu'à la dernière date
                        filtered_vl = fcp_vl.filter(date__gte=ref_vl.date)
                    
                    valeurs = list(filtered_vl.values_list('valeur', flat=True))
                    if len(valeurs) > 1:
                        rendements = [(float(valeurs[i]) / float(valeurs[i-1]) - 1) for i in range(1, len(valeurs))]
                        if rendements and len(rendements) > 1:
                            moyenne = sum(rendements) / len(rendements)
                            variance = sum((r - moyenne) ** 2 for r in rendements) / (len(rendements) - 1)  # Variance non biaisée
                            vol = (variance ** 0.5) * ANNUALIZATION_FACTOR * 100
                        else:
                            vol = 0
                    else:
                        vol = 0
                    
                    # Récupérer le type de fond
                    try:
                        fiche = FicheSignaletique.objects.get(nom=fcp_name)
                        type_fond = fiche.type_fond
                    except FicheSignaletique.DoesNotExist:
                        type_fond = 'Diversifié'
                    
                    all_fcp_stats.append({
                        'nom': fcp_name,
                        'rendement': round(rendement, 2),
                        'volatilite': round(vol, 2),
                        'type_fond': type_fond
                    })
    
    # Récupérer la dernière date de la base (prendre la première disponible)
    last_date_str = None
    for fcp_name in fcp_list:
        fcp_vl_model = get_vl_model(fcp_name)
        if fcp_vl_model:
            last_vl = fcp_vl_model.objects.order_by('-date').first()
            if last_vl:
                last_date_str = last_vl.date.strftime('%d/%m/%Y')
                break
    
    return JsonResponse({'data': all_fcp_stats, 'last_date': last_date_str})


def api_vl_data(request):
    """API pour récupérer les données VL d'un FCP"""
    fcp_name = request.GET.get('fcp')
    period = request.GET.get('period', 'all')
    
    if not fcp_name:
        return JsonResponse({'error': 'FCP non spécifié'}, status=400)
    
    vl_model = get_vl_model(fcp_name)
    if not vl_model:
        return JsonResponse({'error': 'FCP non trouvé'}, status=404)
    
    # Filtrer par période
    vl_queryset = vl_model.objects.all().order_by('date')
    
    if vl_queryset.exists():
        latest_date = vl_queryset.last().date
        
        if period == '1m':
            start_date = latest_date - timedelta(days=30)
        elif period == '3m':
            start_date = latest_date - timedelta(days=90)
        elif period == '6m':
            start_date = latest_date - timedelta(days=180)
        elif period == '1y':
            start_date = latest_date - timedelta(days=365)
        elif period == '5y':
            start_date = latest_date - timedelta(days=365*5)
        else:
            start_date = None
        
        if start_date:
            vl_queryset = vl_queryset.filter(date__gte=start_date)
    
    data = []
    for vl in vl_queryset:
        data.append({
            'date': vl.date.strftime('%Y-%m-%d'),
            'valeur': float(vl.valeur)
        })
    
    return JsonResponse({'data': data})


def api_fcp_full_data(request):
    """API pour récupérer toutes les données d'un FCP (pour mise à jour dynamique)"""
    fcp_name = request.GET.get('fcp')
    
    if not fcp_name:
        return JsonResponse({'error': 'FCP non spécifié'}, status=400)
    
    vl_model = get_vl_model(fcp_name)
    if not vl_model:
        return JsonResponse({'error': 'FCP non trouvé'}, status=404)
    
    vl_data = []
    stats = {}
    perf_calendaires = {}
    perf_glissantes = {}
    analyse_stats = {}
    tracking_error = {}
    
    vl_queryset = vl_model.objects.all().order_by('date')
    
    # Convertir en liste pour JSON
    for vl in vl_queryset:
        vl_data.append({
            'date': vl.date.strftime('%Y-%m-%d'),
            'valeur': float(vl.valeur)
        })
    
    # Calculer les statistiques
    if vl_queryset.exists():
        latest_vl = vl_queryset.last()
        first_vl = vl_queryset.first()
        today = latest_vl.date
        
        def calc_perf(vl_ref):
            if vl_ref:
                return round(((float(latest_vl.valeur) / float(vl_ref.valeur)) - 1) * 100, 2)
            return None
        
        # VL pour différentes périodes GLISSANTES
        vl_1d = vl_queryset.filter(date__lt=today).last()
        vl_1m = vl_queryset.filter(date__lte=today - timedelta(days=30)).last()
        vl_3m = vl_queryset.filter(date__lte=today - timedelta(days=90)).last()
        vl_6m = vl_queryset.filter(date__lte=today - timedelta(days=180)).last()
        vl_1y = vl_queryset.filter(date__lte=today - timedelta(days=365)).last()
        vl_3y = vl_queryset.filter(date__lte=today - timedelta(days=365*3)).last()
        vl_5y = vl_queryset.filter(date__lte=today - timedelta(days=365*5)).last()
        
        # Performances calendaires (To-Date)
        start_of_week = today - timedelta(days=today.weekday())
        vl_wtd = vl_queryset.filter(date__lt=start_of_week).last()
        
        start_of_month = today.replace(day=1)
        vl_mtd = vl_queryset.filter(date__lt=start_of_month).last()
        
        quarter_month = ((today.month - 1) // 3) * 3 + 1
        start_of_quarter = today.replace(month=quarter_month, day=1)
        vl_qtd = vl_queryset.filter(date__lt=start_of_quarter).last()
        
        semester_month = 1 if today.month <= 6 else 7
        start_of_semester = today.replace(month=semester_month, day=1)
        vl_std = vl_queryset.filter(date__lt=start_of_semester).last()
        
        start_of_year = today.replace(month=1, day=1)
        vl_ytd = vl_queryset.filter(date__lt=start_of_year).last()
        
        perf_calendaires = {
            'wtd': calc_perf(vl_wtd),
            'mtd': calc_perf(vl_mtd),
            'qtd': calc_perf(vl_qtd),
            'std': calc_perf(vl_std),
            'ytd': calc_perf(vl_ytd),
        }
        
        perf_glissantes = {
            'perf_1m': calc_perf(vl_1m),
            'perf_3m': calc_perf(vl_3m),
            'perf_6m': calc_perf(vl_6m),
            'perf_1y': calc_perf(vl_1y),
            'perf_3y': calc_perf(vl_3y),
            'perf_5y': calc_perf(vl_5y),
            'origine': calc_perf(first_vl),
        }
        
        stats = {
            'derniere_vl': float(latest_vl.valeur),
            'derniere_date': latest_vl.date.strftime('%d/%m/%Y'),
            'premiere_vl': float(first_vl.valeur),
            'premiere_date': first_vl.date.strftime('%d/%m/%Y'),
            'var_1j': calc_perf(vl_1d) or 0,
            'var_1m': calc_perf(vl_1m) or 0,
            'var_1y': calc_perf(vl_1y) or 0,
            'var_ytd': perf_calendaires['ytd'] or 0,
            'var_origine': perf_glissantes['origine'] or 0,
            'nb_vl': vl_queryset.count(),
        }
        
        # Calculer la Tracking Error pour différentes périodes
        def calc_tracking_error(start_date):
            if start_date is None:
                return None
            period_vl = list(vl_queryset.filter(date__gte=start_date).values_list('valeur', flat=True))
            if len(period_vl) < 2:
                return None
            period_valeurs = [float(v) for v in period_vl]
            period_rendements = [(period_valeurs[i] / period_valeurs[i-1] - 1) * 100 for i in range(1, len(period_valeurs))]
            if len(period_rendements) < 2:
                return None
            moyenne = sum(period_rendements) / len(period_rendements)
            variance = sum((r - moyenne) ** 2 for r in period_rendements) / (len(period_rendements) - 1)  # Variance non biaisée
            ecart_type = variance ** 0.5
            volatilite_ann = ecart_type * ANNUALIZATION_FACTOR
            return round(volatilite_ann, 2)
        
        tracking_error = {
            'wtd': calc_tracking_error(start_of_week),
            'mtd': calc_tracking_error(start_of_month),
            'qtd': calc_tracking_error(start_of_quarter),
            'std': calc_tracking_error(start_of_semester),
            'ytd': calc_tracking_error(start_of_year),
            'origine': calc_tracking_error(first_vl.date if first_vl else None),
        }
        
        # Calculer les statistiques pour l'analyse
        valeurs = [float(v) for v in vl_queryset.values_list('valeur', flat=True)]
        if len(valeurs) > 1:
            rendements = [(valeurs[i] / valeurs[i-1] - 1) * 100 for i in range(1, len(valeurs))]
            
            moyenne_rdt = sum(rendements) / len(rendements)
            variance = sum((r - moyenne_rdt) ** 2 for r in rendements) / (len(rendements) - 1)  # Variance non biaisée
            ecart_type = variance ** 0.5
            volatilite_ann = ecart_type * ANNUALIZATION_FACTOR
            
            rendements_sorted = sorted(rendements)
            n = len(rendements_sorted)
            mediane = rendements_sorted[n // 2] if n % 2 == 1 else (rendements_sorted[n//2 - 1] + rendements_sorted[n//2]) / 2
            
            vl_min = min(valeurs)
            vl_max = max(valeurs)
            rdt_min = min(rendements)
            rdt_max = max(rendements)
            
            jours_positifs = sum(1 for r in rendements if r > 0)
            jours_negatifs = sum(1 for r in rendements if r < 0)
            
            # Drawdown max
            peak = valeurs[0]
            max_drawdown = 0
            for v in valeurs:
                if v > peak:
                    peak = v
                drawdown = (peak - v) / peak * 100
                if drawdown > max_drawdown:
                    max_drawdown = drawdown
            
            sharpe = ((moyenne_rdt - RISK_FREE_RATE_DAILY) / ecart_type * ANNUALIZATION_FACTOR) if ecart_type > 0 else 0
            
            rendements_negatifs = [r for r in rendements if r < 0]
            if rendements_negatifs:
                downside_var = sum(r ** 2 for r in rendements_negatifs) / (len(rendements_negatifs) - 1) if len(rendements_negatifs) > 1 else sum(r ** 2 for r in rendements_negatifs)
                downside_std = downside_var ** 0.5
                sortino = ((moyenne_rdt - RISK_FREE_RATE_DAILY) / downside_std * ANNUALIZATION_FACTOR) if downside_std > 0 else 0
            else:
                sortino = 0
            
            # Skewness et Kurtosis
            if ecart_type > 0:
                skewness = sum((r - moyenne_rdt) ** 3 for r in rendements) / (n * ecart_type ** 3)
                kurtosis = sum((r - moyenne_rdt) ** 4 for r in rendements) / (n * ecart_type ** 4) - 3
            else:
                skewness = 0
                kurtosis = 0
            
            # VaR et CVaR
            var_95 = rendements_sorted[int(n * 0.05)]
            var_99 = rendements_sorted[int(n * 0.01)]
            cvar_95_values = [r for r in rendements if r <= var_95]
            cvar_99_values = [r for r in rendements if r <= var_99]
            cvar_95 = sum(cvar_95_values) / len(cvar_95_values) if cvar_95_values else var_95
            cvar_99 = sum(cvar_99_values) / len(cvar_99_values) if cvar_99_values else var_99
            
            # Drawdowns détaillés
            drawdowns_data = []
            dates_list = list(vl_queryset.values_list('date', flat=True))
            peak = valeurs[0]
            current_dd_start = None
            current_dd_peak = valeurs[0]
            underwater_data = []
            
            for i, v in enumerate(valeurs):
                if v > peak:
                    if current_dd_start is not None:
                        dd_depth = (current_dd_peak - min(valeurs[current_dd_start:i])) / current_dd_peak * 100
                        dd_trough_idx = current_dd_start + valeurs[current_dd_start:i].index(min(valeurs[current_dd_start:i]))
                        drawdowns_data.append({
                            'start_date': dates_list[current_dd_start].strftime('%Y-%m-%d'),
                            'trough_date': dates_list[dd_trough_idx].strftime('%Y-%m-%d'),
                            'end_date': dates_list[i].strftime('%Y-%m-%d'),
                            'depth': round(dd_depth, 2),
                            'duration': i - current_dd_start,
                            'recovery': i - dd_trough_idx
                        })
                        current_dd_start = None
                    peak = v
                    current_dd_peak = v
                else:
                    if current_dd_start is None and v < peak:
                        current_dd_start = i - 1 if i > 0 else 0
                        current_dd_peak = peak
                
                dd_current = (peak - v) / peak * 100
                underwater_data.append({
                    'date': dates_list[i].strftime('%Y-%m-%d'),
                    'drawdown': round(-dd_current, 2)
                })
            
            if current_dd_start is not None:
                dd_depth = (current_dd_peak - min(valeurs[current_dd_start:])) / current_dd_peak * 100
                dd_trough_idx = current_dd_start + valeurs[current_dd_start:].index(min(valeurs[current_dd_start:]))
                drawdowns_data.append({
                    'start_date': dates_list[current_dd_start].strftime('%Y-%m-%d'),
                    'trough_date': dates_list[dd_trough_idx].strftime('%Y-%m-%d'),
                    'end_date': None,
                    'depth': round(dd_depth, 2),
                    'duration': len(valeurs) - current_dd_start,
                    'recovery': None
                })
            
            drawdowns_data.sort(key=lambda x: x['depth'], reverse=True)
            
            # Histogramme
            hist_min = min(rendements)
            hist_max = max(rendements)
            nb_bins = 30
            bin_width = (hist_max - hist_min) / nb_bins if hist_max != hist_min else 1
            histogram_data = []
            for i in range(nb_bins):
                bin_start = hist_min + i * bin_width
                bin_end = bin_start + bin_width
                count = sum(1 for r in rendements if bin_start <= r < bin_end)
                histogram_data.append({
                    'bin_start': round(bin_start, 3),
                    'bin_end': round(bin_end, 3),
                    'count': count,
                    'frequency': round(count / n * 100, 2)
                })
            
            stats['volatilite'] = round(volatilite_ann, 2)
            
            analyse_stats = {
                'nb_observations': len(rendements),
                'rendement_moyen': round(moyenne_rdt, 4),
                'rendement_moyen_ann': round(moyenne_rdt * 365, 2),
                'mediane': round(mediane, 4),
                'ecart_type': round(ecart_type, 4),
                'volatilite_ann': round(volatilite_ann, 2),
                'vl_min': round(vl_min, 2),
                'vl_max': round(vl_max, 2),
                'rdt_min': round(rdt_min, 2),
                'rdt_max': round(rdt_max, 2),
                'max_drawdown': round(max_drawdown, 2),
                'sharpe': round(sharpe, 2),
                'sortino': round(sortino, 2),
                'jours_positifs': jours_positifs,
                'jours_negatifs': jours_negatifs,
                'ratio_positif': round(jours_positifs / len(rendements) * 100, 1) if rendements else 0,
                'skewness': round(skewness, 3),
                'kurtosis': round(kurtosis, 3),
                'var_95': round(var_95, 3),
                'var_99': round(var_99, 3),
                'cvar_95': round(cvar_95, 3),
                'cvar_99': round(cvar_99, 3),
                'histogram_data': histogram_data,
                'underwater_data': underwater_data,
                'drawdowns_data': drawdowns_data[:10],
                'rendements_list': [round(r, 4) for r in rendements],
            }
        else:
            stats['volatilite'] = 0
            tracking_error = {}
    
    # Récupérer les données de la fiche signalétique
    fiche_signaletique_data = None
    try:
        fiche = FicheSignaletique.objects.get(nom=fcp_name)
        fiche_signaletique_data = {
            'type_fond': fiche.type_fond,
            'echelle_risque': fiche.echelle_risque,
            'horizon': fiche.horizon,
            'benchmark_oblig': float(fiche.benchmark_oblig) if fiche.benchmark_oblig else 0,
            'benchmark_brvmc': float(fiche.benchmark_brvmc) if fiche.benchmark_brvmc else 0,
            'description': fiche.description or '',
            'devise': fiche.devise,
            'gestionnaire': fiche.gestionnaire,
            'date_creation': fiche.date_creation.strftime('%d/%m/%Y') if fiche.date_creation else None,
            'depositaire': 'CGF Bourse',
            'frais_gestion': None,  # À ajouter si disponible dans le modèle
            'frais_entree': None,
            'frais_sortie': None,
        }
    except FicheSignaletique.DoesNotExist:
        pass
    
    return JsonResponse({
        'fcp_name': fcp_name,
        'vl_data': vl_data,
        'stats': stats,
        'perf_calendaires': perf_calendaires,
        'perf_glissantes': perf_glissantes,
        'analyse_stats': analyse_stats,
        'tracking_error': tracking_error,
        'fiche_signaletique': fiche_signaletique_data
    })


def composition(request):
    """Vue pour la page de Composition"""
    context = {
        'page_title': 'Composition',
        'page_description': 'Composition détaillée du portefeuille FCP',
        'fcp_list': get_all_fcp_names(),
    }
    return render(request, 'fcp_app/composition.html', context)


# NOTE: Fonction fiche_signaletique supprimée - code orphelin (pas de route ni de template)
# Les données de la fiche signalétique sont intégrées dans la page Valeurs Liquidatives


def a_propos(request):
    """Vue pour la page A Propos"""
    context = {
        'page_title': 'À Propos',
        'page_description': 'Informations sur l\'application de reporting FCP',
        'total_fcp': len(FCP_FICHE_SIGNALETIQUE),
    }
    return render(request, 'fcp_app/a_propos.html', context)


def exportations(request):
    """Vue pour la page des Exportations"""
    # Liste des FCP depuis la base
    fcp_list = list(FicheSignaletique.objects.values_list('nom', flat=True).order_by('nom'))
    
    # Dates par défaut (1 an glissant)
    today = datetime.now().date()
    default_end_date = today.strftime('%Y-%m-%d')
    default_start_date = (today - timedelta(days=365)).strftime('%Y-%m-%d')
    
    context = {
        'page_title': 'Exportations',
        'page_description': 'Exportez vos données FCP en PPT, PDF ou CSV/XLSX',
        'fcp_list': fcp_list,
        'default_start_date': default_start_date,
        'default_end_date': default_end_date,
    }
    return render(request, 'fcp_app/exportations.html', context)


@csrf_exempt
def api_export_data(request):
    """API pour exporter les données FCP en CSV ou XLSX"""
    if request.method != 'POST':
        return JsonResponse({'error': 'Méthode non autorisée'}, status=405)
    
    try:
        data = json.loads(request.body)
        export_format = data.get('format', 'xlsx')
        fcps = data.get('fcps', [])
        content_types = data.get('content', ['vl'])
        start_date_str = data.get('startDate')
        end_date_str = data.get('endDate')
        
        if not fcps:
            return JsonResponse({'error': 'Aucun FCP sélectionné'}, status=400)
        
        # Parser les dates
        start_date = datetime.strptime(start_date_str, '%Y-%m-%d').date() if start_date_str else None
        end_date = datetime.strptime(end_date_str, '%Y-%m-%d').date() if end_date_str else None
        
        # Collecter les données
        all_data = {}
        
        for fcp_name in fcps:
            vl_model = get_vl_model(fcp_name)
            if vl_model:
                queryset = vl_model.objects.all().order_by('date')
                
                if start_date:
                    queryset = queryset.filter(date__gte=start_date)
                if end_date:
                    queryset = queryset.filter(date__lte=end_date)
                
                vl_data = list(queryset.values('date', 'valeur'))
                
                # Calculer les rendements si demandé
                if 'returns' in content_types and len(vl_data) > 1:
                    for i in range(1, len(vl_data)):
                        prev_val = float(vl_data[i-1]['valeur'])
                        curr_val = float(vl_data[i]['valeur'])
                        vl_data[i]['rendement'] = round(((curr_val / prev_val) - 1) * 100, 4) if prev_val else 0
                    vl_data[0]['rendement'] = 0
                
                all_data[fcp_name] = vl_data
        
        if export_format == 'csv':
            return export_to_csv(all_data, content_types)
        else:
            return export_to_xlsx(all_data, content_types)
            
    except json.JSONDecodeError:
        return JsonResponse({'error': 'Données JSON invalides'}, status=400)
    except Exception as e:
        return JsonResponse({'error': str(e)}, status=500)


def export_to_csv(all_data, content_types):
    """Exporter les données en CSV"""
    response = HttpResponse(content_type='text/csv; charset=utf-8')
    response['Content-Disposition'] = f'attachment; filename="export_fcp_{datetime.now().strftime("%Y%m%d_%H%M%S")}.csv"'
    response.write('\ufeff')  # BOM pour Excel
    
    writer = csv.writer(response, delimiter=';')
    
    # En-têtes
    headers = ['FCP', 'Date', 'Valeur Liquidative']
    if 'returns' in content_types:
        headers.append('Rendement (%)')
    writer.writerow(headers)
    
    # Données
    for fcp_name, vl_data in all_data.items():
        for row in vl_data:
            line = [fcp_name, row['date'].strftime('%d/%m/%Y'), str(row['valeur']).replace('.', ',')]
            if 'returns' in content_types:
                rendement = row.get('rendement', 0)
                line.append(str(rendement).replace('.', ','))
            writer.writerow(line)
    
    return response


def export_to_xlsx(all_data, content_types):
    """Exporter les données en XLSX"""
    try:
        import openpyxl
        from openpyxl.styles import Font, Alignment, PatternFill, Border, Side
        from openpyxl.utils import get_column_letter
    except ImportError:
        # Fallback vers CSV si openpyxl n'est pas installé
        return export_to_csv(all_data, content_types)
    
    wb = openpyxl.Workbook()
    
    # Style
    header_font = Font(bold=True, color='FFFFFF')
    header_fill = PatternFill(start_color='004080', end_color='004080', fill_type='solid')
    header_alignment = Alignment(horizontal='center', vertical='center')
    thin_border = Border(
        left=Side(style='thin'),
        right=Side(style='thin'),
        top=Side(style='thin'),
        bottom=Side(style='thin')
    )
    
    # Créer une feuille par FCP
    first_sheet = True
    for fcp_name, vl_data in all_data.items():
        if first_sheet:
            ws = wb.active
            ws.title = fcp_name[:31]  # Max 31 caractères pour le nom de feuille
            first_sheet = False
        else:
            ws = wb.create_sheet(title=fcp_name[:31])
        
        # En-têtes
        headers = ['Date', 'Valeur Liquidative']
        if 'returns' in content_types:
            headers.append('Rendement (%)')
        
        for col, header in enumerate(headers, 1):
            cell = ws.cell(row=1, column=col, value=header)
            cell.font = header_font
            cell.fill = header_fill
            cell.alignment = header_alignment
            cell.border = thin_border
        
        # Données
        for row_idx, row in enumerate(vl_data, 2):
            ws.cell(row=row_idx, column=1, value=row['date']).border = thin_border
            ws.cell(row=row_idx, column=2, value=float(row['valeur'])).border = thin_border
            ws.cell(row=row_idx, column=2).number_format = '#,##0.0000'
            
            if 'returns' in content_types:
                ws.cell(row=row_idx, column=3, value=row.get('rendement', 0)).border = thin_border
                ws.cell(row=row_idx, column=3).number_format = '0.00%'
        
        # Ajuster la largeur des colonnes
        for col in range(1, len(headers) + 1):
            ws.column_dimensions[get_column_letter(col)].width = 18
    
    # Sauvegarder dans un buffer
    buffer = io.BytesIO()
    wb.save(buffer)
    buffer.seek(0)
    
    response = HttpResponse(
        buffer.getvalue(),
        content_type='application/vnd.openxmlformats-officedocument.spreadsheetml.sheet'
    )
    response['Content-Disposition'] = f'attachment; filename="export_fcp_{datetime.now().strftime("%Y%m%d_%H%M%S")}.xlsx"'
    
    return response


@csrf_exempt
def api_export_ppt(request):
    """API pour exporter les données FCP en PowerPoint"""
    if request.method != 'POST':
        return JsonResponse({'error': 'Méthode non autorisée'}, status=405)
    
    try:
        data = json.loads(request.body)
        fcps = data.get('fcps', [])
        content_types = data.get('content', ['vl', 'perf'])
        template = data.get('template', 'standard')
        start_date_str = data.get('startDate')
        end_date_str = data.get('endDate')
        
        if not fcps:
            return JsonResponse({'error': 'Aucun FCP sélectionné'}, status=400)
        
        # Parser les dates
        start_date = datetime.strptime(start_date_str, '%Y-%m-%d').date() if start_date_str else None
        end_date = datetime.strptime(end_date_str, '%Y-%m-%d').date() if end_date_str else None
        
        return generate_ppt(fcps, content_types, template, start_date, end_date)
            
    except json.JSONDecodeError:
        return JsonResponse({'error': 'Données JSON invalides'}, status=400)
    except Exception as e:
        return JsonResponse({'error': str(e)}, status=500)


def generate_ppt(fcps, content_types, template, start_date, end_date):
    """Générer une présentation PowerPoint avec différents templates"""
    from pptx import Presentation
    from pptx.util import Inches, Pt, Emu
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
    from pptx.enum.shapes import MSO_SHAPE
    from pptx.chart.data import CategoryChartData
    from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION
    import matplotlib
    matplotlib.use('Agg')
    import matplotlib.pyplot as plt
    
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    
    # Définition des palettes de couleurs selon le template
    if template == 'moderne':
        # Palette moderne - dégradé violet/bleu
        COLORS = {
            'primary': RGBColor(102, 51, 153),      # Violet
            'secondary': RGBColor(75, 0, 130),      # Indigo
            'accent': RGBColor(138, 43, 226),       # Bleu-violet
            'success': RGBColor(0, 206, 209),       # Turquoise
            'danger': RGBColor(255, 99, 71),        # Tomate
            'light': RGBColor(248, 245, 255),       # Lavande clair
            'dark': RGBColor(45, 45, 60),           # Gris foncé bleuté
            'white': RGBColor(255, 255, 255),
            'chart_line': '#6633cc',
            'chart_fill': '#e8e0f0'
        }
    elif template == 'minimaliste':
        # Palette minimaliste - noir et blanc avec accent
        COLORS = {
            'primary': RGBColor(33, 33, 33),        # Noir
            'secondary': RGBColor(66, 66, 66),      # Gris foncé
            'accent': RGBColor(255, 87, 34),        # Orange accent
            'success': RGBColor(76, 175, 80),       # Vert
            'danger': RGBColor(244, 67, 54),        # Rouge
            'light': RGBColor(250, 250, 250),       # Presque blanc
            'dark': RGBColor(33, 33, 33),           # Noir
            'white': RGBColor(255, 255, 255),
            'chart_line': '#212121',
            'chart_fill': '#f5f5f5'
        }
    else:
        # Palette standard - CGF Bourse
        COLORS = {
            'primary': RGBColor(0, 64, 128),        # Bleu CGF
            'secondary': RGBColor(0, 90, 160),      # Bleu plus clair
            'accent': RGBColor(0, 150, 200),        # Bleu accent
            'success': RGBColor(40, 167, 69),       # Vert
            'danger': RGBColor(220, 53, 69),        # Rouge
            'light': RGBColor(245, 247, 250),       # Gris bleuté clair
            'dark': RGBColor(51, 51, 51),           # Gris foncé
            'white': RGBColor(255, 255, 255),
            'chart_line': '#004080',
            'chart_fill': '#e6f0f7'
        }
    
    def add_title_slide(title, subtitle=""):
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        
        if template == 'minimaliste':
            # Fond blanc avec ligne accent
            line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, Inches(3.5), prs.slide_width, Inches(0.05))
            line.fill.solid()
            line.fill.fore_color.rgb = COLORS['accent']
            line.line.fill.background()
            
            # Titre sobre
            txBox = slide.shapes.add_textbox(Inches(0.5), Inches(2.8), Inches(12.333), Inches(0.8))
            p = txBox.text_frame.paragraphs[0]
            p.text = title.upper()
            p.font.size = Pt(36)
            p.font.bold = True
            p.font.color.rgb = COLORS['dark']
            p.alignment = PP_ALIGN.CENTER
            
            if subtitle:
                txBox = slide.shapes.add_textbox(Inches(0.5), Inches(3.7), Inches(12.333), Inches(0.5))
                p = txBox.text_frame.paragraphs[0]
                p.text = subtitle
                p.font.size = Pt(16)
                p.font.color.rgb = COLORS['secondary']
                p.alignment = PP_ALIGN.CENTER
        
        elif template == 'moderne':
            # Fond dégradé simulé avec formes
            shape1 = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
            shape1.fill.solid()
            shape1.fill.fore_color.rgb = COLORS['primary']
            shape1.line.fill.background()
            
            # Cercles décoratifs
            circle1 = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(-2), Inches(-2), Inches(6), Inches(6))
            circle1.fill.solid()
            circle1.fill.fore_color.rgb = COLORS['secondary']
            circle1.line.fill.background()
            
            circle2 = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(10), Inches(4), Inches(5), Inches(5))
            circle2.fill.solid()
            circle2.fill.fore_color.rgb = COLORS['accent']
            circle2.line.fill.background()
            
            # Titre
            txBox = slide.shapes.add_textbox(Inches(0.5), Inches(2.5), Inches(12.333), Inches(1.2))
            p = txBox.text_frame.paragraphs[0]
            p.text = title
            p.font.size = Pt(48)
            p.font.bold = True
            p.font.color.rgb = COLORS['white']
            p.alignment = PP_ALIGN.CENTER
            
            if subtitle:
                txBox = slide.shapes.add_textbox(Inches(0.5), Inches(4.0), Inches(12.333), Inches(0.6))
                p = txBox.text_frame.paragraphs[0]
                p.text = subtitle
                p.font.size = Pt(20)
                p.font.color.rgb = COLORS['light']
                p.alignment = PP_ALIGN.CENTER
        
        else:  # Standard
            # Fond plein classique
            shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
            shape.fill.solid()
            shape.fill.fore_color.rgb = COLORS['primary']
            shape.line.fill.background()
            
            # Barre décorative
            bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(4), Inches(4.0), Inches(5.333), Inches(0.08))
            bar.fill.solid()
            bar.fill.fore_color.rgb = COLORS['white']
            bar.line.fill.background()
            
            # Titre
            txBox = slide.shapes.add_textbox(Inches(0.5), Inches(2.5), Inches(12.333), Inches(1.2))
            p = txBox.text_frame.paragraphs[0]
            p.text = title
            p.font.size = Pt(44)
            p.font.bold = True
            p.font.color.rgb = COLORS['white']
            p.alignment = PP_ALIGN.CENTER
            
            if subtitle:
                txBox = slide.shapes.add_textbox(Inches(0.5), Inches(4.3), Inches(12.333), Inches(0.6))
                p = txBox.text_frame.paragraphs[0]
                p.text = subtitle
                p.font.size = Pt(22)
                p.font.color.rgb = COLORS['light']
                p.alignment = PP_ALIGN.CENTER
        
        return slide
    
    def add_content_slide(title, fcp_name=None):
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        
        if template == 'minimaliste':
            # Ligne supérieure fine
            line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(0.03))
            line.fill.solid()
            line.fill.fore_color.rgb = COLORS['dark']
            line.line.fill.background()
            
            # Titre aligné à gauche
            txBox = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(8), Inches(0.6))
            p = txBox.text_frame.paragraphs[0]
            p.text = title
            p.font.size = Pt(24)
            p.font.bold = True
            p.font.color.rgb = COLORS['dark']
            
            if fcp_name:
                txBox = slide.shapes.add_textbox(Inches(9), Inches(0.35), Inches(4), Inches(0.5))
                p = txBox.text_frame.paragraphs[0]
                p.text = fcp_name
                p.font.size = Pt(14)
                p.font.color.rgb = COLORS['accent']
                p.alignment = PP_ALIGN.RIGHT
        
        elif template == 'moderne':
            # Bande gauche colorée
            side_bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(0.15), prs.slide_height)
            side_bar.fill.solid()
            side_bar.fill.fore_color.rgb = COLORS['primary']
            side_bar.line.fill.background()
            
            # Header dégradé
            header = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.15), 0, Inches(13.18), Inches(1.0))
            header.fill.solid()
            header.fill.fore_color.rgb = COLORS['light']
            header.line.fill.background()
            
            # Titre
            txBox = slide.shapes.add_textbox(Inches(0.5), Inches(0.25), Inches(8), Inches(0.6))
            p = txBox.text_frame.paragraphs[0]
            p.text = title
            p.font.size = Pt(26)
            p.font.bold = True
            p.font.color.rgb = COLORS['primary']
            
            if fcp_name:
                # Badge FCP
                badge = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(10), Inches(0.2), Inches(3), Inches(0.55))
                badge.fill.solid()
                badge.fill.fore_color.rgb = COLORS['primary']
                badge.line.fill.background()
                
                txBox = slide.shapes.add_textbox(Inches(10), Inches(0.3), Inches(3), Inches(0.4))
                p = txBox.text_frame.paragraphs[0]
                p.text = fcp_name[:25]
                p.font.size = Pt(12)
                p.font.color.rgb = COLORS['white']
                p.alignment = PP_ALIGN.CENTER
        
        else:  # Standard
            # Header classique
            header = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(1.1))
            header.fill.solid()
            header.fill.fore_color.rgb = COLORS['primary']
            header.line.fill.background()
            
            # Titre
            txBox = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.6))
            p = txBox.text_frame.paragraphs[0]
            p.text = title
            p.font.size = Pt(26)
            p.font.bold = True
            p.font.color.rgb = COLORS['white']
            
            if fcp_name:
                txBox = slide.shapes.add_textbox(Inches(9.5), Inches(0.35), Inches(3.5), Inches(0.45))
                p = txBox.text_frame.paragraphs[0]
                p.text = fcp_name
                p.font.size = Pt(13)
                p.font.color.rgb = COLORS['light']
                p.alignment = PP_ALIGN.RIGHT
        
        return slide
    
    def add_kpi_box(slide, label, value, x, y, width=3.8, height=1.4, value_color=None):
        """Ajouter une boîte KPI stylisée selon le template"""
        if value_color is None:
            value_color = COLORS['primary']
        
        if template == 'minimaliste':
            # Style minimal - bordure fine
            box = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(y), Inches(width), Inches(height))
            box.fill.solid()
            box.fill.fore_color.rgb = COLORS['white']
            box.line.color.rgb = COLORS['secondary']
            box.line.width = Pt(1)
        elif template == 'moderne':
            # Style moderne - ombre simulée
            shadow = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x + 0.05), Inches(y + 0.05), Inches(width), Inches(height))
            shadow.fill.solid()
            shadow.fill.fore_color.rgb = RGBColor(200, 200, 210)
            shadow.line.fill.background()
            
            box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(width), Inches(height))
            box.fill.solid()
            box.fill.fore_color.rgb = COLORS['white']
            box.line.fill.background()
        else:
            # Style standard
            box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(width), Inches(height))
            box.fill.solid()
            box.fill.fore_color.rgb = COLORS['light']
            box.line.fill.background()
        
        # Label
        txBox = slide.shapes.add_textbox(Inches(x + 0.15), Inches(y + 0.1), Inches(width - 0.3), Inches(0.35))
        p = txBox.text_frame.paragraphs[0]
        p.text = label
        p.font.size = Pt(11)
        p.font.color.rgb = COLORS['dark']
        
        # Valeur
        txBox = slide.shapes.add_textbox(Inches(x + 0.15), Inches(y + 0.45), Inches(width - 0.3), Inches(0.7))
        p = txBox.text_frame.paragraphs[0]
        p.text = str(value)
        p.font.size = Pt(26)
        p.font.bold = True
        p.font.color.rgb = value_color
    
    def create_chart_image(vl_list, fcp_name):
        """Créer une image de graphique avec matplotlib"""
        fig, ax = plt.subplots(figsize=(11, 4), dpi=100)
        
        dates = [vl.date for vl in vl_list]
        values = [float(vl.valeur) for vl in vl_list]
        
        # Style selon template
        if template == 'moderne':
            ax.fill_between(dates, values, alpha=0.3, color='#6633cc')
            ax.plot(dates, values, color='#6633cc', linewidth=2.5)
            fig.patch.set_facecolor('#f8f5ff')
            ax.set_facecolor('#f8f5ff')
        elif template == 'minimaliste':
            ax.plot(dates, values, color='#212121', linewidth=2)
            fig.patch.set_facecolor('white')
            ax.set_facecolor('white')
            ax.spines['top'].set_visible(False)
            ax.spines['right'].set_visible(False)
        else:
            ax.fill_between(dates, values, alpha=0.2, color='#004080')
            ax.plot(dates, values, color='#004080', linewidth=2)
            fig.patch.set_facecolor('#f5f7fa')
            ax.set_facecolor('#f5f7fa')
        
        ax.set_title(f'Évolution de la VL - {fcp_name}', fontsize=14, fontweight='bold', pad=10)
        ax.set_xlabel('Date', fontsize=10)
        ax.set_ylabel('Valeur Liquidative', fontsize=10)
        ax.grid(True, alpha=0.3, linestyle='--')
        
        # Format des dates
        fig.autofmt_xdate()
        
        plt.tight_layout()
        
        # Sauvegarder en buffer
        img_buffer = io.BytesIO()
        plt.savefig(img_buffer, format='png', bbox_inches='tight', facecolor=fig.get_facecolor())
        img_buffer.seek(0)
        plt.close(fig)
        
        return img_buffer
    
    def add_table(slide, data, headers, x, y, col_widths=None):
        """Ajouter un tableau stylisé"""
        rows = len(data) + 1
        cols = len(headers)
        
        if col_widths is None:
            total_width = 5.5
            col_widths = [total_width / cols] * cols
        
        table_width = sum(col_widths)
        table = slide.shapes.add_table(rows, cols, Inches(x), Inches(y), 
                                        Inches(table_width), Inches(0.38 * rows)).table
        
        # Style selon template
        if template == 'minimaliste':
            header_bg = COLORS['dark']
            alt_bg = COLORS['light']
        elif template == 'moderne':
            header_bg = COLORS['primary']
            alt_bg = RGBColor(248, 245, 255)
        else:
            header_bg = COLORS['primary']
            alt_bg = COLORS['light']
        
        # Headers
        for col_idx, header in enumerate(headers):
            cell = table.cell(0, col_idx)
            cell.text = header
            cell.fill.solid()
            cell.fill.fore_color.rgb = header_bg
            para = cell.text_frame.paragraphs[0]
            para.font.bold = True
            para.font.size = Pt(11)
            para.font.color.rgb = COLORS['white']
            para.alignment = PP_ALIGN.CENTER
        
        # Data
        for row_idx, row_data in enumerate(data, 1):
            for col_idx, value in enumerate(row_data):
                cell = table.cell(row_idx, col_idx)
                cell.text = str(value)
                para = cell.text_frame.paragraphs[0]
                para.font.size = Pt(10)
                para.alignment = PP_ALIGN.CENTER
                
                if row_idx % 2 == 0:
                    cell.fill.solid()
                    cell.fill.fore_color.rgb = alt_bg
        
        return table
    
    def format_perf(value):
        if value is None:
            return "N/A"
        sign = "+" if value >= 0 else ""
        return f"{sign}{value:.2f}%"
    
    # ==================== GÉNÉRATION DES SLIDES ====================
    
    # Slide de titre
    date_range = ""
    if start_date and end_date:
        date_range = f"Période: {start_date.strftime('%d/%m/%Y')} - {end_date.strftime('%d/%m/%Y')}"
    add_title_slide("Rapport FCP", date_range)
    
    # Slide sommaire
    summary_slide = add_content_slide("Sommaire")
    
    y_pos = 1.4 if template == 'minimaliste' else 1.3
    txBox = summary_slide.shapes.add_textbox(Inches(0.5), Inches(y_pos), Inches(12), Inches(0.5))
    p = txBox.text_frame.paragraphs[0]
    p.text = f"Ce rapport analyse {len(fcps)} FCP sur la période sélectionnée"
    p.font.size = Pt(14)
    p.font.color.rgb = COLORS['dark']
    
    # Liste des FCP en colonnes
    col1_x, col2_x = 0.8, 6.8
    for i, fcp in enumerate(fcps):
        x = col1_x if i % 2 == 0 else col2_x
        y = 2.0 + (i // 2) * 0.5
        
        txBox = summary_slide.shapes.add_textbox(Inches(x), Inches(y), Inches(5.5), Inches(0.4))
        p = txBox.text_frame.paragraphs[0]
        p.text = f"● {fcp}"
        p.font.size = Pt(13)
        p.font.color.rgb = COLORS['primary'] if template != 'minimaliste' else COLORS['dark']
    
    # Slides par FCP
    for fcp_name in fcps:
        vl_model = get_vl_model(fcp_name)
        if not vl_model:
            continue
        
        queryset = vl_model.objects.all().order_by('date')
        if start_date:
            queryset = queryset.filter(date__gte=start_date)
        if end_date:
            queryset = queryset.filter(date__lte=end_date)
        
        vl_list = list(queryset)
        if not vl_list:
            continue
        
        latest_vl = vl_list[-1]
        first_vl = vl_list[0]
        perf_periode = ((float(latest_vl.valeur) / float(first_vl.valeur)) - 1) * 100
        
        # ===== SLIDE VL avec graphique =====
        if 'vl' in content_types:
            vl_slide = add_content_slide("Évolution des Valeurs Liquidatives", fcp_name)
            
            # KPIs en haut
            kpi_y = 1.25
            add_kpi_box(vl_slide, "Dernière VL", f"{float(latest_vl.valeur):,.4f}", 0.4, kpi_y, 3.0, 1.2)
            add_kpi_box(vl_slide, "Date", latest_vl.date.strftime('%d/%m/%Y'), 3.6, kpi_y, 2.4, 1.2)
            
            perf_color = COLORS['success'] if perf_periode >= 0 else COLORS['danger']
            add_kpi_box(vl_slide, "Performance", format_perf(perf_periode), 6.2, kpi_y, 2.6, 1.2, perf_color)
            add_kpi_box(vl_slide, "Observations", str(len(vl_list)), 9.0, kpi_y, 2.4, 1.2)
            
            # Min/Max
            min_vl = min(vl_list, key=lambda x: float(x.valeur))
            max_vl = max(vl_list, key=lambda x: float(x.valeur))
            add_kpi_box(vl_slide, "VL Min", f"{float(min_vl.valeur):,.4f}", 11.6, kpi_y, 1.5, 1.2)
            
            # GRAPHIQUE DE COURBE
            if len(vl_list) >= 2:
                chart_img = create_chart_image(vl_list, fcp_name)
                vl_slide.shapes.add_picture(chart_img, Inches(0.4), Inches(2.7), width=Inches(12.5), height=Inches(4.5))
        
        # ===== SLIDE PERFORMANCES =====
        if 'perf' in content_types:
            perf_slide = add_content_slide("Analyse des Performances", fcp_name)
            
            all_vl = vl_model.objects.all().order_by('date')
            latest_date = latest_vl.date
            
            def get_perf(days_back=None, start_of_period=None):
                if start_of_period:
                    ref = all_vl.filter(date__lt=start_of_period).last()
                elif days_back:
                    ref = all_vl.filter(date__lte=latest_date - timedelta(days=days_back)).last()
                else:
                    return None
                if ref:
                    return ((float(latest_vl.valeur) / float(ref.valeur)) - 1) * 100
                return None
            
            # Calculs
            start_year = latest_date.replace(month=1, day=1)
            start_month = latest_date.replace(day=1)
            
            ytd = get_perf(start_of_period=start_year)
            mtd = get_perf(start_of_period=start_month)
            perf_1m = get_perf(days_back=30)
            perf_3m = get_perf(days_back=90)
            perf_6m = get_perf(days_back=180)
            perf_1y = get_perf(days_back=365)
            perf_3y = get_perf(days_back=1095)
            
            # KPIs principaux en haut
            kpi_y = 1.3
            mtd_color = COLORS['success'] if mtd and mtd >= 0 else COLORS['danger']
            ytd_color = COLORS['success'] if ytd and ytd >= 0 else COLORS['danger']
            
            add_kpi_box(perf_slide, "MTD", format_perf(mtd), 0.5, kpi_y, 3.0, 1.3, mtd_color if mtd else COLORS['dark'])
            add_kpi_box(perf_slide, "YTD", format_perf(ytd), 3.7, kpi_y, 3.0, 1.3, ytd_color if ytd else COLORS['dark'])
            add_kpi_box(perf_slide, "1 An", format_perf(perf_1y), 6.9, kpi_y, 3.0, 1.3)
            add_kpi_box(perf_slide, "3 Ans", format_perf(perf_3y), 10.1, kpi_y, 2.9, 1.3)
            
            # Tableaux de performances
            # Calendaires
            cal_data = [
                ["MTD", format_perf(mtd)],
                ["YTD", format_perf(ytd)],
                ["Période", format_perf(perf_periode)],
            ]
            
            txBox = perf_slide.shapes.add_textbox(Inches(0.5), Inches(2.9), Inches(5), Inches(0.4))
            p = txBox.text_frame.paragraphs[0]
            p.text = "Performances Calendaires"
            p.font.size = Pt(14)
            p.font.bold = True
            p.font.color.rgb = COLORS['primary'] if template != 'minimaliste' else COLORS['dark']
            
            add_table(perf_slide, cal_data, ["Période", "Perf."], 0.5, 3.4, [2.5, 2.5])
            
            # Glissantes
            gliss_data = [
                ["1 Mois", format_perf(perf_1m)],
                ["3 Mois", format_perf(perf_3m)],
                ["6 Mois", format_perf(perf_6m)],
                ["1 An", format_perf(perf_1y)],
                ["3 Ans", format_perf(perf_3y)],
            ]
            
            txBox = perf_slide.shapes.add_textbox(Inches(6.5), Inches(2.9), Inches(5), Inches(0.4))
            p = txBox.text_frame.paragraphs[0]
            p.text = "Performances Glissantes"
            p.font.size = Pt(14)
            p.font.bold = True
            p.font.color.rgb = COLORS['primary'] if template != 'minimaliste' else COLORS['dark']
            
            add_table(perf_slide, gliss_data, ["Horizon", "Perf."], 6.5, 3.4, [2.5, 2.5])
            
            # Statistiques
            values = [float(vl.valeur) for vl in vl_list]
            avg_vl = sum(values) / len(values)
            
            txBox = perf_slide.shapes.add_textbox(Inches(0.5), Inches(5.8), Inches(12), Inches(0.4))
            p = txBox.text_frame.paragraphs[0]
            p.text = f"VL Moyenne: {avg_vl:,.4f}  |  VL Min: {min(values):,.4f}  |  VL Max: {max(values):,.4f}  |  Observations: {len(values)}"
            p.font.size = Pt(11)
            p.font.color.rgb = COLORS['dark']
        
        # ===== SLIDE FICHE SIGNALÉTIQUE =====
        if 'fiche' in content_types:
            try:
                fiche = FicheSignaletique.objects.get(nom=fcp_name)
                fiche_slide = add_content_slide("Fiche Signalétique", fcp_name)
                
                fiche_data = [
                    ["Type de fond", fiche.type_fond or "N/A"],
                    ["Échelle de risque", f"{fiche.echelle_risque}/7 - {fiche.risk_label}"],
                    ["Horizon", f"{fiche.horizon} ans" if fiche.horizon else "N/A"],
                    ["Date création", fiche.date_creation.strftime('%d/%m/%Y') if fiche.date_creation else "N/A"],
                    ["Devise", fiche.devise or "XOF"],
                    ["Gestionnaire", fiche.gestionnaire or "N/A"],
                ]
                
                # Tableau centré
                add_table(fiche_slide, fiche_data, ["Caractéristique", "Valeur"], 3.5, 1.8, [3.0, 4.0])
                
                # Description si disponible
                if fiche.description:
                    txBox = fiche_slide.shapes.add_textbox(Inches(1), Inches(5.0), Inches(11), Inches(1.5))
                    tf = txBox.text_frame
                    tf.word_wrap = True
                    p = tf.paragraphs[0]
                    p.text = fiche.description[:300]
                    p.font.size = Pt(11)
                    p.font.color.rgb = COLORS['dark']
                
            except FicheSignaletique.DoesNotExist:
                pass
    
    # Slide de fin
    add_title_slide("Merci", f"Rapport généré le {datetime.now().strftime('%d/%m/%Y à %H:%M')}")
    
    # Sauvegarder
    buffer = io.BytesIO()
    prs.save(buffer)
    buffer.seek(0)
    
    response = HttpResponse(
        buffer.getvalue(),
        content_type='application/vnd.openxmlformats-officedocument.presentationml.presentation'
    )
    response['Content-Disposition'] = f'attachment; filename="rapport_fcp_{template}_{datetime.now().strftime("%Y%m%d_%H%M%S")}.pptx"'
    
    return response


@csrf_exempt
def api_export_pdf(request):
    """API pour exporter les données FCP en PDF"""
    if request.method != 'POST':
        return JsonResponse({'error': 'Méthode non autorisée'}, status=405)
    
    try:
        data = json.loads(request.body)
        fcps = data.get('fcps', [])
        content_types = data.get('content', ['vl', 'perf'])
        start_date_str = data.get('startDate')
        end_date_str = data.get('endDate')
        
        if not fcps:
            return JsonResponse({'error': 'Aucun FCP sélectionné'}, status=400)
        
        # Parser les dates
        start_date = datetime.strptime(start_date_str, '%Y-%m-%d').date() if start_date_str else None
        end_date = datetime.strptime(end_date_str, '%Y-%m-%d').date() if end_date_str else None
        
        return generate_pdf(fcps, content_types, start_date, end_date)
            
    except json.JSONDecodeError:
        return JsonResponse({'error': 'Données JSON invalides'}, status=400)
    except Exception as e:
        return JsonResponse({'error': str(e)}, status=500)


def generate_pdf(fcps, content_types, start_date, end_date):
    """Générer un rapport PDF"""
    from reportlab.lib import colors
    from reportlab.lib.pagesizes import A4, landscape
    from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
    from reportlab.lib.units import cm, mm
    from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer, Table, TableStyle, PageBreak
    from reportlab.lib.enums import TA_CENTER, TA_LEFT, TA_RIGHT
    
    buffer = io.BytesIO()
    doc = SimpleDocTemplate(buffer, pagesize=A4, 
                           rightMargin=1.5*cm, leftMargin=1.5*cm,
                           topMargin=2*cm, bottomMargin=2*cm)
    
    # Styles
    styles = getSampleStyleSheet()
    
    title_style = ParagraphStyle(
        'CustomTitle',
        parent=styles['Heading1'],
        fontSize=24,
        textColor=colors.HexColor('#004080'),
        spaceAfter=30,
        alignment=TA_CENTER
    )
    
    heading_style = ParagraphStyle(
        'CustomHeading',
        parent=styles['Heading2'],
        fontSize=16,
        textColor=colors.HexColor('#004080'),
        spaceBefore=20,
        spaceAfter=10
    )
    
    subheading_style = ParagraphStyle(
        'CustomSubHeading',
        parent=styles['Heading3'],
        fontSize=14,
        textColor=colors.HexColor('#333333'),
        spaceBefore=15,
        spaceAfter=8
    )
    
    normal_style = ParagraphStyle(
        'CustomNormal',
        parent=styles['Normal'],
        fontSize=10,
        textColor=colors.HexColor('#333333'),
        spaceAfter=6
    )
    
    # Contenu
    story = []
    
    # Page de titre
    story.append(Spacer(1, 3*cm))
    story.append(Paragraph("Rapport FCP", title_style))
    
    date_range = f"{start_date.strftime('%d/%m/%Y')} - {end_date.strftime('%d/%m/%Y')}" if start_date and end_date else ""
    story.append(Paragraph(f"Période: {date_range}", normal_style))
    story.append(Paragraph(f"Généré le {datetime.now().strftime('%d/%m/%Y à %H:%M')}", normal_style))
    story.append(Spacer(1, 2*cm))
    
    # Liste des FCP
    story.append(Paragraph("FCP inclus dans ce rapport:", subheading_style))
    for fcp in fcps:
        story.append(Paragraph(f"• {fcp}", normal_style))
    
    story.append(PageBreak())
    
    # Contenu par FCP
    for fcp_name in fcps:
        vl_model = get_vl_model(fcp_name)
        if not vl_model:
            continue
        
        queryset = vl_model.objects.all().order_by('date')
        if start_date:
            queryset = queryset.filter(date__gte=start_date)
        if end_date:
            queryset = queryset.filter(date__lte=end_date)
        
        vl_list = list(queryset)
        if not vl_list:
            continue
        
        latest_vl = vl_list[-1] if vl_list else None
        first_vl = vl_list[0] if vl_list else None
        
        # Titre FCP
        story.append(Paragraph(fcp_name, heading_style))
        
        # VL
        if 'vl' in content_types and latest_vl:
            story.append(Paragraph("Valeurs Liquidatives", subheading_style))
            
            # Stats
            perf_periode = ((float(latest_vl.valeur) / float(first_vl.valeur)) - 1) * 100 if first_vl else 0
            
            stats_data = [
                ["Dernière VL", f"{float(latest_vl.valeur):,.4f}", latest_vl.date.strftime('%d/%m/%Y')],
                ["Première VL (période)", f"{float(first_vl.valeur):,.4f}", first_vl.date.strftime('%d/%m/%Y')],
                ["Performance période", f"{'+' if perf_periode >= 0 else ''}{perf_periode:.2f}%", ""],
                ["Nombre de points", str(len(vl_list)), ""],
            ]
            
            stats_table = Table(stats_data, colWidths=[5*cm, 4*cm, 4*cm])
            stats_table.setStyle(TableStyle([
                ('BACKGROUND', (0, 0), (-1, -1), colors.HexColor('#f5f5f4')),
                ('TEXTCOLOR', (0, 0), (-1, -1), colors.HexColor('#333333')),
                ('ALIGN', (0, 0), (-1, -1), 'LEFT'),
                ('FONTNAME', (0, 0), (-1, -1), 'Helvetica'),
                ('FONTSIZE', (0, 0), (-1, -1), 10),
                ('BOTTOMPADDING', (0, 0), (-1, -1), 8),
                ('TOPPADDING', (0, 0), (-1, -1), 8),
                ('GRID', (0, 0), (-1, -1), 0.5, colors.HexColor('#e0dedd')),
            ]))
            story.append(stats_table)
            story.append(Spacer(1, 0.5*cm))
            
            # Tableau dernières VL
            story.append(Paragraph("Dernières valeurs liquidatives:", normal_style))
            recent_vls = vl_list[-10:] if len(vl_list) > 10 else vl_list
            
            vl_data = [["Date", "Valeur Liquidative"]]
            for vl in reversed(recent_vls):
                vl_data.append([vl.date.strftime('%d/%m/%Y'), f"{float(vl.valeur):,.4f}"])
            
            vl_table = Table(vl_data, colWidths=[5*cm, 5*cm])
            vl_table.setStyle(TableStyle([
                ('BACKGROUND', (0, 0), (-1, 0), colors.HexColor('#004080')),
                ('TEXTCOLOR', (0, 0), (-1, 0), colors.white),
                ('ALIGN', (0, 0), (-1, -1), 'CENTER'),
                ('FONTNAME', (0, 0), (-1, 0), 'Helvetica-Bold'),
                ('FONTSIZE', (0, 0), (-1, -1), 9),
                ('BOTTOMPADDING', (0, 0), (-1, -1), 6),
                ('TOPPADDING', (0, 0), (-1, -1), 6),
                ('GRID', (0, 0), (-1, -1), 0.5, colors.HexColor('#e0dedd')),
                ('ROWBACKGROUNDS', (0, 1), (-1, -1), [colors.white, colors.HexColor('#f5f5f4')]),
            ]))
            story.append(vl_table)
            story.append(Spacer(1, 0.5*cm))
        
        # Performances
        if 'perf' in content_types and latest_vl:
            story.append(Paragraph("Performances", subheading_style))
            
            all_vl = vl_model.objects.all().order_by('date')
            latest_date = latest_vl.date
            
            def get_perf(days_back=None, start_of_period=None):
                if start_of_period:
                    ref = all_vl.filter(date__lt=start_of_period).last()
                elif days_back:
                    ref = all_vl.filter(date__lte=latest_date - timedelta(days=days_back)).last()
                else:
                    return None
                if ref:
                    return ((float(latest_vl.valeur) / float(ref.valeur)) - 1) * 100
                return None
            
            def format_perf(val):
                if val is None:
                    return "N/A"
                return f"{'+' if val >= 0 else ''}{val:.2f}%"
            
            # YTD & MTD
            start_year = latest_date.replace(month=1, day=1)
            start_month = latest_date.replace(day=1)
            
            perf_data = [
                ["Période", "Performance"],
                ["MTD", format_perf(get_perf(start_of_period=start_month))],
                ["YTD", format_perf(get_perf(start_of_period=start_year))],
                ["1 Mois", format_perf(get_perf(days_back=30))],
                ["3 Mois", format_perf(get_perf(days_back=90))],
                ["6 Mois", format_perf(get_perf(days_back=180))],
                ["1 An", format_perf(get_perf(days_back=365))],
            ]
            
            perf_table = Table(perf_data, colWidths=[5*cm, 5*cm])
            perf_table.setStyle(TableStyle([
                ('BACKGROUND', (0, 0), (-1, 0), colors.HexColor('#004080')),
                ('TEXTCOLOR', (0, 0), (-1, 0), colors.white),
                ('ALIGN', (0, 0), (-1, -1), 'CENTER'),
                ('FONTNAME', (0, 0), (-1, 0), 'Helvetica-Bold'),
                ('FONTSIZE', (0, 0), (-1, -1), 10),
                ('BOTTOMPADDING', (0, 0), (-1, -1), 8),
                ('TOPPADDING', (0, 0), (-1, -1), 8),
                ('GRID', (0, 0), (-1, -1), 0.5, colors.HexColor('#e0dedd')),
                ('ROWBACKGROUNDS', (0, 1), (-1, -1), [colors.white, colors.HexColor('#f5f5f4')]),
            ]))
            story.append(perf_table)
            story.append(Spacer(1, 0.5*cm))
        
        # Fiche signalétique
        if 'fiche' in content_types:
            try:
                fiche = FicheSignaletique.objects.get(nom=fcp_name)
                story.append(Paragraph("Fiche Signalétique", subheading_style))
                
                fiche_data = [
                    ["Caractéristique", "Valeur"],
                    ["Type de fond", fiche.type_fond or "N/A"],
                    ["Échelle de risque", f"{fiche.echelle_risque}/7 - {fiche.risk_label}"],
                    ["Horizon d'investissement", f"{fiche.horizon} ans" if fiche.horizon else "N/A"],
                    ["Date de création", fiche.date_creation.strftime('%d/%m/%Y') if fiche.date_creation else "N/A"],
                    ["Devise", fiche.devise or "XOF"],
                    ["Gestionnaire", fiche.gestionnaire or "N/A"],
                ]
                
                fiche_table = Table(fiche_data, colWidths=[5*cm, 8*cm])
                fiche_table.setStyle(TableStyle([
                    ('BACKGROUND', (0, 0), (-1, 0), colors.HexColor('#004080')),
                    ('TEXTCOLOR', (0, 0), (-1, 0), colors.white),
                    ('ALIGN', (0, 0), (0, -1), 'LEFT'),
                    ('ALIGN', (1, 0), (1, -1), 'LEFT'),
                    ('FONTNAME', (0, 0), (-1, 0), 'Helvetica-Bold'),
                    ('FONTSIZE', (0, 0), (-1, -1), 10),
                    ('BOTTOMPADDING', (0, 0), (-1, -1), 8),
                    ('TOPPADDING', (0, 0), (-1, -1), 8),
                    ('GRID', (0, 0), (-1, -1), 0.5, colors.HexColor('#e0dedd')),
                    ('ROWBACKGROUNDS', (0, 1), (-1, -1), [colors.white, colors.HexColor('#f5f5f4')]),
                ]))
                story.append(fiche_table)
            except FicheSignaletique.DoesNotExist:
                pass
        
        story.append(PageBreak())
    
    # Build PDF
    doc.build(story)
    buffer.seek(0)
    
    response = HttpResponse(buffer.getvalue(), content_type='application/pdf')
    response['Content-Disposition'] = f'attachment; filename="rapport_fcp_{datetime.now().strftime("%Y%m%d_%H%M%S")}.pdf"'
    
    return response


def api_factsheet_preview(request):
    """API pour récupérer les données de prévisualisation du factsheet"""
    try:
        fcp_name = request.GET.get('fcp')
        month = request.GET.get('month')  # Format: YYYY-MM
        
        if not fcp_name:
            return JsonResponse({'error': 'FCP non spécifié'}, status=400)
        
        if not month:
            today = datetime.now().date()
            if today.month == 1:
                month = f"{today.year - 1}-12"
            else:
                month = f"{today.year}-{today.month - 1:02d}"
        
        # Obtenir le modèle VL du FCP
        vl_model = get_vl_model(fcp_name)
        if not vl_model:
            return JsonResponse({'error': f'FCP "{fcp_name}" non trouvé'}, status=404)
        
        # Récupérer la fiche signalétique
        try:
            fiche = FicheSignaletique.objects.get(nom_fcp=fcp_name)
            fiche_data = {
                'type_fond': fiche.type_fond or 'N/A',
                'gestionnaire': fiche.gestionnaire or 'N/A',
                'devise': fiche.devise or 'XOF',
                'echelle_risque': fiche.echelle_risque or 'N/A',
                'risk_label': fiche.risk_label or 'N/A',
                'horizon': fiche.horizon or 'N/A',
                'date_creation': fiche.date_creation.strftime('%d/%m/%Y') if fiche.date_creation else 'N/A',
            }
        except FicheSignaletique.DoesNotExist:
            fiche_data = {
                'type_fond': 'N/A',
                'gestionnaire': 'CGF Bourse',
                'devise': 'XOF',
                'echelle_risque': 'N/A',
                'risk_label': 'N/A',
                'horizon': 'N/A',
                'date_creation': 'N/A',
            }
        
        # Parsing du mois
        year, month_num = int(month.split('-')[0]), int(month.split('-')[1])
        end_of_month = datetime(year, month_num, 1) + timedelta(days=32)
        end_of_month = datetime(end_of_month.year, end_of_month.month, 1) - timedelta(days=1)
        
        # Récupérer toutes les VL jusqu'à la fin du mois
        vl_queryset = vl_model.objects.filter(date__lte=end_of_month.date()).order_by('date')
        vl_list = list(vl_queryset.values('date', 'valeur'))
        
        if not vl_list:
            return JsonResponse({'error': 'Aucune valeur liquidative disponible'}, status=404)
        
        latest_vl = vl_list[-1]
        
        # Fonctions de calcul de performance
        def get_perf(days_back=None, target_date=None):
            if target_date:
                filtered = [v for v in vl_list if v['date'] <= target_date]
                if len(filtered) < 2:
                    return None
                return (float(filtered[-1]['valeur']) / float(filtered[0]['valeur']) - 1) * 100
            
            if days_back is None or len(vl_list) < 2:
                return None
            
            ref_date = latest_vl['date'] - timedelta(days=days_back)
            for v in reversed(vl_list[:-1]):
                if v['date'] <= ref_date:
                    return (float(latest_vl['valeur']) / float(v['valeur']) - 1) * 100
            return None
        
        def get_perf_since(start_date):
            filtered = [v for v in vl_list if v['date'] >= start_date]
            if len(filtered) < 2:
                return None
            return (float(filtered[-1]['valeur']) / float(filtered[0]['valeur']) - 1) * 100
        
        # Calculs des performances
        perf_1m = get_perf(30)
        perf_3m = get_perf(91)
        perf_6m = get_perf(182)
        perf_1y = get_perf(365)
        perf_3y = get_perf(365 * 3)
        perf_5y = get_perf(365 * 5)
        perf_origine = (float(latest_vl['valeur']) / float(vl_list[0]['valeur']) - 1) * 100 if vl_list else None
        
        # Performances To-Date
        today = latest_vl['date']
        # WTD - depuis lundi
        days_since_monday = today.weekday()
        start_of_week = today - timedelta(days=days_since_monday)
        perf_wtd = get_perf_since(start_of_week)
        
        # MTD - depuis début du mois
        start_of_month = today.replace(day=1)
        perf_mtd = get_perf_since(start_of_month)
        
        # QTD - depuis début du trimestre
        quarter_month = ((today.month - 1) // 3) * 3 + 1
        start_of_quarter = today.replace(month=quarter_month, day=1)
        perf_qtd = get_perf_since(start_of_quarter)
        
        # STD - depuis début du semestre
        semester_month = 1 if today.month <= 6 else 7
        start_of_semester = today.replace(month=semester_month, day=1)
        perf_std = get_perf_since(start_of_semester)
        
        # YTD - depuis début de l'année
        start_of_year = today.replace(month=1, day=1)
        perf_ytd = get_perf_since(start_of_year)
        
        # Calcul des statistiques
        volatilite = None
        sharpe = None
        max_dd = None
        
        if len(vl_list) > 1:
            valeurs = [float(v['valeur']) for v in vl_list]
            rendements = [(valeurs[i] / valeurs[i-1] - 1) * 100 for i in range(1, len(valeurs))]
            
            moyenne_rdt = sum(rendements) / len(rendements)
            variance = sum((r - moyenne_rdt) ** 2 for r in rendements) / (len(rendements) - 1) if len(rendements) > 1 else sum((r - moyenne_rdt) ** 2 for r in rendements)
            volatilite = (variance ** 0.5) * ANNUALIZATION_FACTOR
            
            sharpe = ((moyenne_rdt - RISK_FREE_RATE_DAILY) / (variance ** 0.5) * ANNUALIZATION_FACTOR) if variance > 0 else 0
            
            peak = valeurs[0]
            max_dd = 0
            for v in valeurs:
                if v > peak:
                    peak = v
                dd = (peak - v) / peak * 100
                if dd > max_dd:
                    max_dd = dd
        
        # Formater les performances
        def fmt_perf(val):
            if val is None:
                return 'N/A'
            sign = '+' if val >= 0 else ''
            return f'{sign}{val:.2f}%'
        
        # Préparer les données pour le chart (base 100)
        chart_data = []
        if vl_list:
            first_val = float(vl_list[0]['valeur'])
            # Échantillonner max 50 points
            step = max(1, len(vl_list) // 50)
            sampled = vl_list[::step]
            if vl_list[-1] not in sampled:
                sampled.append(vl_list[-1])
            for vl in sampled:
                chart_data.append({
                    'date': vl['date'].strftime('%Y-%m-%d'),
                    'value': round(float(vl['valeur']) / first_val * 100, 2)
                })
        
        # Mois formaté
        month_names = ['Janvier', 'Février', 'Mars', 'Avril', 'Mai', 'Juin', 
                       'Juillet', 'Août', 'Septembre', 'Octobre', 'Novembre', 'Décembre']
        formatted_month = f"{month_names[month_num - 1]} {year}"
        
        return JsonResponse({
            'fcp_name': fcp_name,
            'month': formatted_month,
            'fiche': fiche_data,
            'latest_vl': {
                'date': latest_vl['date'].strftime('%d/%m/%Y'),
                'valeur': f"{float(latest_vl['valeur']):,.4f}".replace(',', ' ')
            },
            'performances': {
                'todate': {
                    'wtd': fmt_perf(perf_wtd),
                    'mtd': fmt_perf(perf_mtd),
                    'qtd': fmt_perf(perf_qtd),
                    'std': fmt_perf(perf_std),
                    'ytd': fmt_perf(perf_ytd),
                },
                'glissantes': {
                    '1m': fmt_perf(perf_1m),
                    '3m': fmt_perf(perf_3m),
                    '6m': fmt_perf(perf_6m),
                    '1y': fmt_perf(perf_1y),
                    '3y': fmt_perf(perf_3y),
                    '5y': fmt_perf(perf_5y),
                    'origine': fmt_perf(perf_origine),
                }
            },
            'statistics': {
                'volatilite': f"{volatilite:.2f}%" if volatilite else 'N/A',
                'sharpe': f"{sharpe:.2f}" if sharpe else 'N/A',
                'max_drawdown': f"-{max_dd:.2f}%" if max_dd else 'N/A',
                'nb_observations': len(vl_list),
            },
            'chart_data': chart_data,
        })
    except Exception as e:
        import traceback
        traceback.print_exc()
        return JsonResponse({'error': str(e)}, status=500)


@csrf_exempt
def api_export_factsheet(request):
    """API pour exporter un Factsheet PDF d'une page pour un FCP"""
    if request.method != 'POST':
        return JsonResponse({'error': 'Méthode non autorisée'}, status=405)
    
    try:
        data = json.loads(request.body)
        fcp_name = data.get('fcp')
        month = data.get('month')  # Format: YYYY-MM
        commentaire = data.get('commentaire', '')
        disclaimer = data.get('disclaimer', "Ce document est fourni à titre informatif uniquement et ne constitue pas une offre ou une sollicitation d'achat ou de vente. Les performances passées ne préjugent pas des performances futures.")
        
        if not fcp_name:
            return JsonResponse({'error': 'FCP non spécifié'}, status=400)
        
        if not month:
            # Par défaut, le mois précédent
            today = datetime.now().date()
            if today.month == 1:
                month = f"{today.year - 1}-12"
            else:
                month = f"{today.year}-{today.month - 1:02d}"
        
        return generate_factsheet_pdf(fcp_name, month, commentaire, disclaimer)
            
    except json.JSONDecodeError:
        return JsonResponse({'error': 'Données JSON invalides'}, status=400)
    except Exception as e:
        import traceback
        traceback.print_exc()
        return JsonResponse({'error': str(e)}, status=500)


def generate_factsheet_pdf(fcp_name, month, commentaire, disclaimer):
    """Générer un Factsheet PDF d'une page A4 paysage"""
    from reportlab.lib import colors
    from reportlab.lib.pagesizes import A4, landscape
    from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
    from reportlab.lib.units import cm, mm
    from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer, Table, TableStyle, Image, KeepTogether, BaseDocTemplate, Frame, PageTemplate
    from reportlab.lib.enums import TA_CENTER, TA_LEFT, TA_RIGHT, TA_JUSTIFY
    from reportlab.graphics.shapes import Drawing, Rect, String, Line, Circle
    from reportlab.graphics.charts.lineplots import LinePlot
    from reportlab.graphics.charts.piecharts import Pie
    from reportlab.graphics.charts.barcharts import VerticalBarChart
    from reportlab.pdfbase import pdfmetrics
    from reportlab.pdfbase.ttfonts import TTFont
    import os
    
    # Essayer d'enregistrer la police Aptos Narrow
    FONT_NAME = 'Helvetica'
    FONT_NAME_BOLD = 'Helvetica-Bold'
    try:
        # Chemins possibles pour Aptos Narrow
        font_paths = [
            'C:/Windows/Fonts/aptos-narrow.ttf',
            'C:/Windows/Fonts/AptosNarrow.ttf',
            'C:/Windows/Fonts/Aptos-Narrow.ttf',
            os.path.expanduser('~/AppData/Local/Microsoft/Windows/Fonts/aptos-narrow.ttf'),
        ]
        for font_path in font_paths:
            if os.path.exists(font_path):
                pdfmetrics.registerFont(TTFont('AptosNarrow', font_path))
                FONT_NAME = 'AptosNarrow'
                break
        # Chercher aussi la version bold
        font_bold_paths = [
            'C:/Windows/Fonts/aptos-narrow-bold.ttf',
            'C:/Windows/Fonts/AptosNarrow-Bold.ttf',
        ]
        for font_path in font_bold_paths:
            if os.path.exists(font_path):
                pdfmetrics.registerFont(TTFont('AptosNarrow-Bold', font_path))
                FONT_NAME_BOLD = 'AptosNarrow-Bold'
                break
    except Exception:
        pass  # Utiliser Helvetica par défaut
    
    # Récupérer les informations du FCP
    try:
        fiche = FicheSignaletique.objects.get(nom=fcp_name)
    except FicheSignaletique.DoesNotExist:
        return JsonResponse({'error': 'FCP non trouvé'}, status=404)
    
    vl_model = get_vl_model(fcp_name)
    if not vl_model:
        return JsonResponse({'error': 'Modèle VL non trouvé'}, status=404)
    
    # Parser le mois
    year, month_num = map(int, month.split('-'))
    mois_fr = {1: 'Janvier', 2: 'Février', 3: 'Mars', 4: 'Avril', 5: 'Mai', 6: 'Juin',
               7: 'Juillet', 8: 'Août', 9: 'Septembre', 10: 'Octobre', 11: 'Novembre', 12: 'Décembre'}
    
    # Récupérer les VL
    vl_queryset = vl_model.objects.all().order_by('date')
    vl_list = list(vl_queryset.values('date', 'valeur'))
    
    if not vl_list:
        return JsonResponse({'error': 'Aucune donnée VL disponible'}, status=400)
    
    # Filtrer pour le mois sélectionné
    end_of_month = datetime(year, month_num + 1, 1).date() if month_num < 12 else datetime(year + 1, 1, 1).date()
    month_vls = [v for v in vl_list if v['date'] < end_of_month]
    if not month_vls:
        month_vls = vl_list
    
    latest_vl = month_vls[-1]
    first_vl = vl_list[0]
    latest_date = latest_vl['date']
    
    # Fonctions utilitaires
    def get_vl_at_date(target_date):
        for vl in reversed(vl_list):
            if vl['date'] <= target_date:
                return float(vl['valeur'])
        return None
    
    def calc_perf_from_vl(ref_val):
        if ref_val:
            return ((float(latest_vl['valeur']) / ref_val) - 1) * 100
        return None
    
    def fmt_perf(val):
        if val is None:
            return "N/A"
        sign = "+" if val >= 0 else ""
        return f"{sign}{val:.2f}%"
    
    def fmt_val(val, suffix=""):
        if val is None:
            return "N/A"
        return f"{val:.2f}{suffix}"
    
    # Calculs des performances
    start_of_year = latest_date.replace(month=1, day=1)
    vl_ytd = get_vl_at_date(start_of_year - timedelta(days=1))
    vl_1y = get_vl_at_date(latest_date - timedelta(days=365))
    vl_3y = get_vl_at_date(latest_date - timedelta(days=365*3))
    vl_5y = get_vl_at_date(latest_date - timedelta(days=365*5))
    vl_origine = float(first_vl['valeur'])
    
    perf_ytd = calc_perf_from_vl(vl_ytd)
    perf_1y = calc_perf_from_vl(vl_1y)
    perf_3y = calc_perf_from_vl(vl_3y)
    perf_5y = calc_perf_from_vl(vl_5y)
    perf_origine = calc_perf_from_vl(vl_origine)
    
    # Performances To-Date
    start_of_week = latest_date - timedelta(days=latest_date.weekday())
    start_of_month_perf = latest_date.replace(day=1)
    quarter_month = ((latest_date.month - 1) // 3) * 3 + 1
    start_of_quarter = latest_date.replace(month=quarter_month, day=1)
    semester_month = 1 if latest_date.month <= 6 else 7
    start_of_semester = latest_date.replace(month=semester_month, day=1)
    
    perf_wtd = calc_perf_from_vl(get_vl_at_date(start_of_week - timedelta(days=1)))
    perf_mtd = calc_perf_from_vl(get_vl_at_date(start_of_month_perf - timedelta(days=1)))
    perf_qtd = calc_perf_from_vl(get_vl_at_date(start_of_quarter - timedelta(days=1)))
    perf_std = calc_perf_from_vl(get_vl_at_date(start_of_semester - timedelta(days=1)))
    
    # Performances glissantes
    perf_1m = calc_perf_from_vl(get_vl_at_date(latest_date - timedelta(days=30)))
    perf_3m = calc_perf_from_vl(get_vl_at_date(latest_date - timedelta(days=90)))
    perf_6m = calc_perf_from_vl(get_vl_at_date(latest_date - timedelta(days=180)))
    
    # Calculer les statistiques de risque
    valeurs = [float(v['valeur']) for v in vl_list]
    rendements = [(valeurs[i] / valeurs[i-1] - 1) * 100 for i in range(1, len(valeurs))]
    
    def calc_stats_for_period(start_date):
        """Calcule les stats pour une période donnée"""
        period_vls = [v for v in vl_list if v['date'] >= start_date]
        if len(period_vls) < 2:
            return {'vol': None, 'max_dd': None, 'var95': None, 'te': None, 'sharpe': None, 'sortino': None}
        
        vals = [float(v['valeur']) for v in period_vls]
        rets = [(vals[i] / vals[i-1] - 1) * 100 for i in range(1, len(vals))]
        
        if not rets:
            return {'vol': None, 'max_dd': None, 'var95': None, 'te': None, 'sharpe': None, 'sortino': None}
        
        mean_ret = sum(rets) / len(rets)
        variance = sum((r - mean_ret) ** 2 for r in rets) / (len(rets) - 1) if len(rets) > 1 else sum((r - mean_ret) ** 2 for r in rets)  # Variance non biaisée
        vol = (variance ** 0.5) * ANNUALIZATION_FACTOR if variance > 0 else 0
        
        # Max Drawdown
        peak = vals[0]
        max_dd = 0
        for v in vals:
            if v > peak:
                peak = v
            dd = (peak - v) / peak * 100
            if dd > max_dd:
                max_dd = dd
        
        # VaR 95%
        sorted_rets = sorted(rets)
        var95_idx = int(len(sorted_rets) * 0.05)
        var95 = sorted_rets[var95_idx] if var95_idx < len(sorted_rets) else sorted_rets[0]
        
        # Tracking Error (volatilité)
        te = vol
        
        # Sharpe (utiliser constantes)
        sharpe = ((mean_ret - RISK_FREE_RATE_DAILY) / (variance ** 0.5) * ANNUALIZATION_FACTOR) if variance > 0 else 0
        
        # Sortino
        neg_rets = [r for r in rets if r < 0]
        if neg_rets:
            downside_var = sum(r ** 2 for r in neg_rets) / (len(neg_rets) - 1) if len(neg_rets) > 1 else sum(r ** 2 for r in neg_rets)
            sortino = ((mean_ret - RISK_FREE_RATE_DAILY) / (downside_var ** 0.5) * ANNUALIZATION_FACTOR) if downside_var > 0 else 0
        else:
            sortino = 0
        
        return {'vol': vol, 'max_dd': max_dd, 'var95': abs(var95), 'te': te, 'sharpe': sharpe, 'sortino': sortino}
    
    stats_ytd = calc_stats_for_period(start_of_year)
    stats_1y = calc_stats_for_period(latest_date - timedelta(days=365))
    stats_3y = calc_stats_for_period(latest_date - timedelta(days=365*3))
    stats_5y = calc_stats_for_period(latest_date - timedelta(days=365*5))
    stats_origine = calc_stats_for_period(first_vl['date'])
    
    # Couleurs
    GRIS = colors.HexColor('#6c757d')
    BLEU = colors.HexColor('#004080')
    GRIS_CLAIR = colors.HexColor('#f8f9fa')
    BLEU_CLAIR = colors.HexColor('#e6f0fa')
    VERT = colors.HexColor('#28a745')
    ROUGE = colors.HexColor('#dc3545')
    
    # Créer le PDF
    buffer = io.BytesIO()
    page_width, page_height = landscape(A4)
    
    # Marges très réduites pour tout faire tenir
    LEFT_MARGIN = 0.5*cm
    RIGHT_MARGIN = 0.5*cm
    TOP_MARGIN = 0.3*cm
    BOTTOM_MARGIN = 0.3*cm
    
    doc = SimpleDocTemplate(
        buffer, 
        pagesize=landscape(A4),
        rightMargin=RIGHT_MARGIN, 
        leftMargin=LEFT_MARGIN,
        topMargin=TOP_MARGIN, 
        bottomMargin=BOTTOM_MARGIN
    )
    
    usable_width = page_width - LEFT_MARGIN - RIGHT_MARGIN
    
    # Styles avec police compacte
    styles = getSampleStyleSheet()
    
    style_header = ParagraphStyle('Header', fontName=FONT_NAME_BOLD, fontSize=9, textColor=colors.white, alignment=TA_LEFT)
    style_header_center = ParagraphStyle('HeaderCenter', fontName=FONT_NAME, fontSize=8, textColor=colors.white, alignment=TA_CENTER)
    style_header_right = ParagraphStyle('HeaderRight', fontName=FONT_NAME, fontSize=8, textColor=colors.white, alignment=TA_RIGHT)
    style_section = ParagraphStyle('Section', fontName=FONT_NAME_BOLD, fontSize=6, textColor=BLEU, spaceBefore=2, spaceAfter=1)
    style_tiny = ParagraphStyle('Tiny', fontName=FONT_NAME, fontSize=5, textColor=colors.HexColor('#333'), leading=6)
    style_disclaimer = ParagraphStyle('Disclaimer', fontName=FONT_NAME, fontSize=4.5, textColor=colors.HexColor('#666'), alignment=TA_JUSTIFY, leading=5.5)
    
    story = []
    
    # ============== EN-TÊTE (barre tricolore) ==============
    header_data = [[
        Paragraph(fcp_name, style_header),
        Paragraph(f"Fonds {fiche.type_fond}", style_header_center),
        Paragraph(f"Rapport mensuel {mois_fr[month_num]} {year}", style_header_right)
    ]]
    header_table = Table(header_data, colWidths=[usable_width*0.38, usable_width*0.24, usable_width*0.38])
    header_table.setStyle(TableStyle([
        ('BACKGROUND', (0, 0), (0, 0), GRIS),
        ('BACKGROUND', (1, 0), (1, 0), BLEU),
        ('BACKGROUND', (2, 0), (2, 0), GRIS),
        ('VALIGN', (0, 0), (-1, -1), 'MIDDLE'),
        ('LEFTPADDING', (0, 0), (-1, -1), 6),
        ('RIGHTPADDING', (0, 0), (-1, -1), 6),
        ('TOPPADDING', (0, 0), (-1, -1), 5),
        ('BOTTOMPADDING', (0, 0), (-1, -1), 5),
    ]))
    story.append(header_table)
    story.append(Spacer(1, 3))
    
    # ============== CORPS PRINCIPAL ==============
    # Largeurs des colonnes: 55% gauche, 45% droite
    LEFT_COL_WIDTH = usable_width * 0.54
    RIGHT_COL_WIDTH = usable_width * 0.46
    
    # ---- COLONNE GAUCHE ----
    def create_perf_chart():
        """Graphique performance base 100 compact"""
        drawing = Drawing(LEFT_COL_WIDTH - 10, 75)
        
        base_100_data = []
        first_val = float(vl_list[0]['valeur'])
        step = max(1, len(vl_list) // 80)
        sampled = vl_list[::step]
        if vl_list[-1] not in sampled:
            sampled.append(vl_list[-1])
        
        for i, vl in enumerate(sampled):
            base_100_data.append(float(vl['valeur']) / first_val * 100)
        
        lp = LinePlot()
        lp.x = 25
        lp.y = 10
        lp.height = 55
        lp.width = LEFT_COL_WIDTH - 50
        lp.data = [[(i, v) for i, v in enumerate(base_100_data)]]
        lp.lines[0].strokeColor = BLEU
        lp.lines[0].strokeWidth = 1
        lp.xValueAxis.visibleLabels = 0
        lp.xValueAxis.strokeColor = colors.HexColor('#ddd')
        lp.yValueAxis.strokeColor = colors.HexColor('#ddd')
        lp.yValueAxis.labelTextFormat = '%d'
        lp.yValueAxis.labels.fontSize = 5
        lp.yValueAxis.labels.fontName = FONT_NAME
        drawing.add(lp)
        drawing.add(String(LEFT_COL_WIDTH/2 - 5, 68, "Performance (base 100)", fontSize=5, fontName=FONT_NAME_BOLD, textAnchor='middle', fillColor=BLEU))
        return drawing
    
    def create_compact_table(title, headers, rows, col_widths):
        """Crée un tableau compact"""
        data = [headers] + rows
        table = Table(data, colWidths=col_widths)
        table.setStyle(TableStyle([
            ('BACKGROUND', (0, 0), (-1, 0), BLEU),
            ('TEXTCOLOR', (0, 0), (-1, 0), colors.white),
            ('BACKGROUND', (0, 1), (0, -1), GRIS_CLAIR),
            ('FONTNAME', (0, 0), (-1, -1), FONT_NAME),
            ('FONTSIZE', (0, 0), (-1, -1), 5),
            ('ALIGN', (0, 0), (-1, -1), 'CENTER'),
            ('VALIGN', (0, 0), (-1, -1), 'MIDDLE'),
            ('GRID', (0, 0), (-1, -1), 0.3, colors.HexColor('#ccc')),
            ('TOPPADDING', (0, 0), (-1, -1), 1.5),
            ('BOTTOMPADDING', (0, 0), (-1, -1), 1.5),
            ('LEFTPADDING', (0, 0), (-1, -1), 2),
            ('RIGHTPADDING', (0, 0), (-1, -1), 2),
        ]))
        return table
    
    # Perf To-Date
    perf_td_headers = ['', 'WTD', 'MTD', 'QTD', 'STD', 'YTD']
    perf_td_rows = [
        ['FCP', fmt_perf(perf_wtd), fmt_perf(perf_mtd), fmt_perf(perf_qtd), fmt_perf(perf_std), fmt_perf(perf_ytd)],
        ['Bench.', 'N/A', 'N/A', 'N/A', 'N/A', 'N/A'],
        ['Diff.', 'N/A', 'N/A', 'N/A', 'N/A', 'N/A'],
    ]
    td_col_w = (LEFT_COL_WIDTH - 10) / 6
    perf_td_table = create_compact_table("Perf. To-Date", perf_td_headers, perf_td_rows, [td_col_w]*6)
    
    # Perf Glissantes
    perf_gl_headers = ['', '1M', '3M', '6M', '1A', '3A', '5A', 'Orig.']
    perf_gl_rows = [
        ['FCP', fmt_perf(perf_1m), fmt_perf(perf_3m), fmt_perf(perf_6m), fmt_perf(perf_1y), fmt_perf(perf_3y), fmt_perf(perf_5y), fmt_perf(perf_origine)],
        ['Bench.', 'N/A', 'N/A', 'N/A', 'N/A', 'N/A', 'N/A', 'N/A'],
        ['Diff.', 'N/A', 'N/A', 'N/A', 'N/A', 'N/A', 'N/A', 'N/A'],
    ]
    gl_col_w = (LEFT_COL_WIDTH - 10) / 8
    perf_gl_table = create_compact_table("Perf. Glissantes", perf_gl_headers, perf_gl_rows, [gl_col_w]*8)
    
    # Graphiques Allocation et Maturité côte à côte
    def create_mini_pie():
        drawing = Drawing((LEFT_COL_WIDTH-10)/2, 55)
        allocations = {'Oblig.': 45, 'Actions': 30, 'Monet.': 15, 'Autres': 10}
        pie = Pie()
        pie.x = 25
        pie.y = 5
        pie.width = 40
        pie.height = 40
        pie.data = list(allocations.values())
        pie.labels = list(allocations.keys())
        pie.slices.strokeWidth = 0.3
        pie.slices.strokeColor = colors.white
        pie.slices.fontName = FONT_NAME
        pie.slices.fontSize = 4
        for i, c in enumerate([BLEU, VERT, colors.HexColor('#ffc107'), GRIS]):
            pie.slices[i].fillColor = c
        pie.sideLabels = 1
        drawing.add(pie)
        drawing.add(String((LEFT_COL_WIDTH-10)/4, 50, "Allocation", fontSize=5, fontName=FONT_NAME_BOLD, textAnchor='middle', fillColor=BLEU))
        return drawing
    
    def create_mini_bars():
        drawing = Drawing((LEFT_COL_WIDTH-10)/2, 55)
        mats = {'<1A': 15, '1-3A': 25, '3-5A': 30, '5-7A': 20, '>7A': 10}
        bc = VerticalBarChart()
        bc.x = 20
        bc.y = 8
        bc.height = 35
        bc.width = (LEFT_COL_WIDTH-10)/2 - 35
        bc.data = [list(mats.values())]
        bc.categoryAxis.categoryNames = list(mats.keys())
        bc.categoryAxis.labels.fontSize = 4
        bc.categoryAxis.labels.fontName = FONT_NAME
        bc.valueAxis.valueMin = 0
        bc.valueAxis.valueMax = 40
        bc.valueAxis.labels.fontSize = 4
        bc.valueAxis.labelTextFormat = '%d%%'
        bc.bars[0].fillColor = BLEU
        bc.barWidth = 8
        drawing.add(bc)
        drawing.add(String((LEFT_COL_WIDTH-10)/4, 50, "Maturité", fontSize=5, fontName=FONT_NAME_BOLD, textAnchor='middle', fillColor=BLEU))
        return drawing
    
    charts_mini = Table([[create_mini_pie(), create_mini_bars()]], colWidths=[(LEFT_COL_WIDTH-10)/2, (LEFT_COL_WIDTH-10)/2])
    
    # Tableau Actions
    act_headers = ['', 'P/E', 'P/B', 'Div Yield', 'Beta']
    act_rows = [['FCP', 'N/A', 'N/A', 'N/A', 'N/A'], ['Bench.', 'N/A', 'N/A', 'N/A', 'N/A'], ['Diff.', 'N/A', 'N/A', 'N/A', 'N/A']]
    act_col_w = (LEFT_COL_WIDTH - 10) / 5
    actions_table = create_compact_table("Actions", act_headers, act_rows, [act_col_w]*5)
    
    # Tableau Obligations
    obl_headers = ['', 'Nb Lig.', 'Duration', 'Sensib.', 'Matur.', 'YTM', 'Coupon']
    obl_rows = [['FCP', 'N/A', 'N/A', 'N/A', 'N/A', 'N/A', 'N/A']]
    obl_col_w = (LEFT_COL_WIDTH - 10) / 7
    oblig_table = create_compact_table("Obligations", obl_headers, obl_rows, [obl_col_w]*7)
    
    # Assembler colonne gauche
    left_elements = [
        [Paragraph('<b>PERFORMANCE</b>', style_section)],
        [create_perf_chart()],
        [Paragraph('<b>Performances To-Date</b>', style_tiny)],
        [perf_td_table],
        [Spacer(1, 2)],
        [Paragraph('<b>Performances Glissantes</b>', style_tiny)],
        [perf_gl_table],
        [Spacer(1, 2)],
        [Paragraph('<b>CARACTÉRISTIQUES DU PORTEFEUILLE</b>', style_section)],
        [charts_mini],
        [Spacer(1, 2)],
        [Paragraph('<b>Sous-ptf Actions</b>', style_tiny)],
        [actions_table],
        [Spacer(1, 2)],
        [Paragraph('<b>Sous-ptf Obligations</b>', style_tiny)],
        [oblig_table],
    ]
    left_table = Table(left_elements, colWidths=[LEFT_COL_WIDTH - 5])
    
    # ---- COLONNE DROITE ----
    # Échelle de risque graphique
    def create_risk_scale():
        drawing = Drawing(RIGHT_COL_WIDTH - 15, 25)
        risk_colors = ['#22c55e', '#84cc16', '#eab308', '#f97316', '#ef4444', '#dc2626', '#991b1b']
        box_w = (RIGHT_COL_WIDTH - 30) / 7
        for i in range(7):
            opacity = 0.3 if i + 1 != fiche.echelle_risque else 1.0
            rect = Rect(5 + i * box_w, 8, box_w - 2, 12)
            rect.fillColor = colors.HexColor(risk_colors[i])
            rect.fillOpacity = opacity
            rect.strokeWidth = 0
            drawing.add(rect)
            drawing.add(String(5 + i * box_w + box_w/2, 11, str(i+1), fontSize=5, fontName=FONT_NAME_BOLD, textAnchor='middle', fillColor=colors.white if opacity == 1 else colors.HexColor('#666')))
        drawing.add(String(5, 3, "Risque faible", fontSize=4, fontName=FONT_NAME, fillColor=colors.HexColor('#666')))
        drawing.add(String(RIGHT_COL_WIDTH - 20, 3, "Risque élevé", fontSize=4, fontName=FONT_NAME, textAnchor='end', fillColor=colors.HexColor('#666')))
        return drawing
    
    # Infos clés
    vl_str = f"{float(latest_vl['valeur']):,.2f}".replace(',', ' ')
    infos_data = [
        [Paragraph(f"<b>VL:</b> {vl_str} XOF au {latest_vl['date'].strftime('%d/%m/%Y')}", style_tiny)],
        [Paragraph(f"<b>Type:</b> {fiche.type_fond} | <b>Risque:</b> {fiche.echelle_risque}/7 | <b>Horizon:</b> {fiche.horizon} ans", style_tiny)],
        [Paragraph(f"<b>Devise:</b> {fiche.devise} | <b>Gestionnaire:</b> {fiche.gestionnaire}", style_tiny)],
    ]
    infos_table = Table(infos_data, colWidths=[RIGHT_COL_WIDTH - 15])
    infos_table.setStyle(TableStyle([
        ('BACKGROUND', (0, 0), (-1, -1), BLEU_CLAIR),
        ('TOPPADDING', (0, 0), (-1, -1), 2),
        ('BOTTOMPADDING', (0, 0), (-1, -1), 2),
        ('LEFTPADDING', (0, 0), (-1, -1), 3),
    ]))
    
    # Tableau Indicateurs de risque
    risk_headers = ['', 'YTD', '1A', '3A', '5A', 'Orig.']
    risk_rows = [
        ['MaxDD', fmt_val(stats_ytd['max_dd'], '%'), fmt_val(stats_1y['max_dd'], '%'), fmt_val(stats_3y['max_dd'], '%'), fmt_val(stats_5y['max_dd'], '%'), fmt_val(stats_origine['max_dd'], '%')],
        ['Track.Err', fmt_val(stats_ytd['te'], '%'), fmt_val(stats_1y['te'], '%'), fmt_val(stats_3y['te'], '%'), fmt_val(stats_5y['te'], '%'), fmt_val(stats_origine['te'], '%')],
        ['Volatilité', fmt_val(stats_ytd['vol'], '%'), fmt_val(stats_1y['vol'], '%'), fmt_val(stats_3y['vol'], '%'), fmt_val(stats_5y['vol'], '%'), fmt_val(stats_origine['vol'], '%')],
        ['VaR 95%', fmt_val(stats_ytd['var95'], '%'), fmt_val(stats_1y['var95'], '%'), fmt_val(stats_3y['var95'], '%'), fmt_val(stats_5y['var95'], '%'), fmt_val(stats_origine['var95'], '%')],
    ]
    risk_col_w = (RIGHT_COL_WIDTH - 15) / 6
    risk_table = create_compact_table("Risque", risk_headers, risk_rows, [risk_col_w]*6)
    
    # Tableau Rendement ajusté au risque
    adj_headers = ['', '1A', '3A', 'Orig.']
    adj_rows = [
        ['Beta', 'N/A', 'N/A', 'N/A'],
        ['Info Ratio', 'N/A', 'N/A', 'N/A'],
        ['Sharpe', fmt_val(stats_1y['sharpe']), fmt_val(stats_3y['sharpe']), fmt_val(stats_origine['sharpe'])],
        ['Sortino', fmt_val(stats_1y['sortino']), fmt_val(stats_3y['sortino']), fmt_val(stats_origine['sortino'])],
    ]
    adj_col_w = (RIGHT_COL_WIDTH - 15) / 4
    adj_table = create_compact_table("Rdt ajusté risque", adj_headers, adj_rows, [adj_col_w]*4)
    
    # Tableau Contribution à la performance
    contrib_headers = ['', 'Ptf', 'Bench.', 'Diff.']
    contrib_rows = [
        ['Actions', 'N/A', 'N/A', 'N/A'],
        ['Obligations', 'N/A', 'N/A', 'N/A'],
        ['OPCVM', 'N/A', 'N/A', 'N/A'],
        ['Trésorerie', 'N/A', 'N/A', 'N/A'],
        ['Frais gestion', 'N/A', 'N/A', 'N/A'],
        ['Frais transac.', 'N/A', 'N/A', 'N/A'],
        ['Total', 'N/A', 'N/A', 'N/A'],
    ]
    contrib_col_w = (RIGHT_COL_WIDTH - 15) / 4
    contrib_table = create_compact_table("Contribution Perf.", contrib_headers, contrib_rows, [contrib_col_w]*4)
    
    # Tableau Attribution à la performance
    attrib_headers = ['', 'Alloc.', 'Sélec.', 'Inter.', 'Total']
    attrib_rows = [
        ['Actions', 'N/A', 'N/A', 'N/A', 'N/A'],
        ['Obligations', 'N/A', 'N/A', 'N/A', 'N/A'],
        ['OPCVM', 'N/A', 'N/A', 'N/A', 'N/A'],
        ['Trésorerie', 'N/A', 'N/A', 'N/A', 'N/A'],
        ['Total', 'N/A', 'N/A', 'N/A', 'N/A'],
    ]
    attrib_col_w = (RIGHT_COL_WIDTH - 15) / 5
    attrib_table = create_compact_table("Attribution Perf.", attrib_headers, attrib_rows, [attrib_col_w]*5)
    
    # Assembler colonne droite
    right_elements = [
        [Paragraph('<b>FICHE SIGNALÉTIQUE</b>', style_section)],
        [create_risk_scale()],
        [infos_table],
        [Spacer(1, 2)],
        [Paragraph('<b>INDICATEURS DE RISQUE</b>', style_section)],
        [risk_table],
        [Spacer(1, 2)],
        [Paragraph('<b>RENDEMENT AJUSTÉ AU RISQUE</b>', style_section)],
        [adj_table],
        [Spacer(1, 2)],
        [Paragraph('<b>CONTRIBUTION À LA PERFORMANCE</b>', style_section)],
        [contrib_table],
        [Spacer(1, 2)],
        [Paragraph('<b>ATTRIBUTION DE LA PERFORMANCE</b>', style_section)],
        [attrib_table],
    ]
    right_table = Table(right_elements, colWidths=[RIGHT_COL_WIDTH - 5])
    
    # Assembler les deux colonnes
    main_table = Table([[left_table, right_table]], colWidths=[LEFT_COL_WIDTH, RIGHT_COL_WIDTH])
    main_table.setStyle(TableStyle([('VALIGN', (0, 0), (-1, -1), 'TOP')]))
    story.append(main_table)
    
    # ============== FOOTER ==============
    story.append(Spacer(1, 3))
    
    # Ligne séparatrice
    sep = Table([['']],  colWidths=[usable_width])
    sep.setStyle(TableStyle([('LINEABOVE', (0, 0), (-1, 0), 0.5, BLEU)]))
    story.append(sep)
    
    # Commentaire et Disclaimer côte à côte pour gagner de la place
    footer_left = ""
    if commentaire:
        footer_left = f"<b>COMMENTAIRE:</b> {commentaire}"
    footer_right = f"<b>AVERTISSEMENT:</b> {disclaimer}"
    
    if commentaire:
        footer_data = [[Paragraph(footer_left, style_disclaimer), Paragraph(footer_right, style_disclaimer)]]
        footer_table = Table(footer_data, colWidths=[usable_width * 0.4, usable_width * 0.6])
    else:
        footer_data = [[Paragraph(footer_right, style_disclaimer)]]
        footer_table = Table(footer_data, colWidths=[usable_width])
    
    footer_table.setStyle(TableStyle([
        ('VALIGN', (0, 0), (-1, -1), 'TOP'),
        ('TOPPADDING', (0, 0), (-1, -1), 2),
    ]))
    story.append(footer_table)
    
    # Build PDF
    doc.build(story)
    buffer.seek(0)
    
    filename = f"factsheet_{fcp_name.replace(' ', '_')}_{month}.pdf"
    response = HttpResponse(buffer.getvalue(), content_type='application/pdf')
    response['Content-Disposition'] = f'attachment; filename="{filename}"'
    
    return response


def api_correlation_matrix(request):
    """API pour calculer la matrice de corrélation entre les FCP"""
    period = request.GET.get('period', 'origin')
    
    fcp_list = list(FicheSignaletique.objects.values_list('nom', flat=True).order_by('nom'))
    
    # Récupérer les rendements de chaque FCP pour la période sélectionnée
    fcp_returns = {}
    fcp_dates = {}
    
    for fcp_name in fcp_list:
        fcp_vl_model = get_vl_model(fcp_name)
        if fcp_vl_model:
            fcp_vl = fcp_vl_model.objects.all().order_by('date')
            if fcp_vl.exists() and fcp_vl.count() > 10:
                last = fcp_vl.last()
                latest_date = last.date
                
                # Déterminer la date de début selon la période
                if period == 'wtd':
                    start_of_period = latest_date - timedelta(days=latest_date.weekday())
                    ref_vl = fcp_vl.filter(date__lt=start_of_period).last()
                elif period == 'mtd':
                    start_of_period = latest_date.replace(day=1)
                    ref_vl = fcp_vl.filter(date__lt=start_of_period).last()
                elif period == 'qtd':
                    quarter_month = ((latest_date.month - 1) // 3) * 3 + 1
                    start_of_period = latest_date.replace(month=quarter_month, day=1)
                    ref_vl = fcp_vl.filter(date__lt=start_of_period).last()
                elif period == 'std':
                    semester_month = 1 if latest_date.month <= 6 else 7
                    start_of_period = latest_date.replace(month=semester_month, day=1)
                    ref_vl = fcp_vl.filter(date__lt=start_of_period).last()
                elif period == 'ytd':
                    start_of_period = latest_date.replace(month=1, day=1)
                    ref_vl = fcp_vl.filter(date__lt=start_of_period).last()
                else:  # origin
                    ref_vl = fcp_vl.first()
                    start_of_period = ref_vl.date if ref_vl else None
                
                if ref_vl:
                    # Filtrer les VL depuis le début de la période
                    filtered_vl = fcp_vl.filter(date__gte=ref_vl.date)
                    
                    # Calculer les rendements quotidiens
                    vl_list = list(filtered_vl.values('date', 'valeur'))
                    if len(vl_list) > 1:
                        returns = {}
                        for i in range(1, len(vl_list)):
                            date_str = vl_list[i]['date'].strftime('%Y-%m-%d')
                            ret = (float(vl_list[i]['valeur']) / float(vl_list[i-1]['valeur']) - 1) * 100
                            returns[date_str] = ret
                        fcp_returns[fcp_name] = returns
                        fcp_dates[fcp_name] = set(returns.keys())
    
    # Trouver les dates communes à tous les FCP
    if fcp_dates:
        common_dates = set.intersection(*fcp_dates.values()) if len(fcp_dates) > 1 else set()
        common_dates = sorted(common_dates)
    else:
        common_dates = []
    
    # Calculer la matrice de corrélation
    fcp_names = list(fcp_returns.keys())
    n = len(fcp_names)
    correlation_matrix = []
    
    if len(common_dates) > 5:  # Au moins 5 observations communes
        for i, fcp1 in enumerate(fcp_names):
            row = []
            returns1 = [fcp_returns[fcp1].get(d, 0) for d in common_dates]
            mean1 = sum(returns1) / len(returns1)
            std1 = (sum((r - mean1) ** 2 for r in returns1) / len(returns1)) ** 0.5
            
            for j, fcp2 in enumerate(fcp_names):
                if i == j:
                    row.append(1.0)
                else:
                    returns2 = [fcp_returns[fcp2].get(d, 0) for d in common_dates]
                    mean2 = sum(returns2) / len(returns2)
                    std2 = (sum((r - mean2) ** 2 for r in returns2) / len(returns2)) ** 0.5
                    
                    if std1 > 0 and std2 > 0:
                        covariance = sum((returns1[k] - mean1) * (returns2[k] - mean2) for k in range(len(common_dates))) / len(common_dates)
                        corr = covariance / (std1 * std2)
                        row.append(round(corr, 3))
                    else:
                        row.append(0)
            correlation_matrix.append(row)
    
    return JsonResponse({
        'fcp_names': fcp_names,
        'matrix': correlation_matrix,
        'nb_observations': len(common_dates),
        'period': period
    })


def api_volatility_clustering(request):
    """API pour le clustering de volatilité"""
    fcp_name = request.GET.get('fcp')
    window = int(request.GET.get('window', 20))  # Fenêtre glissante par défaut 20 jours
    
    # Limiter la fenêtre entre 5 et 30
    window = max(5, min(30, window))
    
    if not fcp_name:
        return JsonResponse({'error': 'FCP non spécifié'}, status=400)
    
    vl_model = get_vl_model(fcp_name)
    if not vl_model:
        return JsonResponse({'error': 'FCP non trouvé'}, status=404)
    
    vl_queryset = vl_model.objects.all().order_by('date')
    
    if not vl_queryset.exists() or vl_queryset.count() < window + 10:
        return JsonResponse({'error': 'Données insuffisantes'}, status=400)
    
    # Calculer les rendements quotidiens
    vl_list = list(vl_queryset.values('date', 'valeur'))
    rendements = []
    dates = []
    
    for i in range(1, len(vl_list)):
        ret = (float(vl_list[i]['valeur']) / float(vl_list[i-1]['valeur']) - 1) * 100
        rendements.append(ret)
        dates.append(vl_list[i]['date'].strftime('%Y-%m-%d'))
    
    # Calculer la volatilité glissante (annualisée)
    volatilites = []
    vol_dates = []
    
    for i in range(window - 1, len(rendements)):
        window_returns = rendements[i - window + 1:i + 1]
        mean_ret = sum(window_returns) / len(window_returns)
        variance = sum((r - mean_ret) ** 2 for r in window_returns) / (len(window_returns) - 1) if len(window_returns) > 1 else sum((r - mean_ret) ** 2 for r in window_returns)
        vol = (variance ** 0.5) * ANNUALIZATION_FACTOR  # Volatilité annualisée
        volatilites.append(vol)
        vol_dates.append(dates[i])
    
    if len(volatilites) < 10:
        return JsonResponse({'error': 'Données insuffisantes pour le clustering'}, status=400)
    
    # Clustering K-means simplifié (3 clusters)
    # Trier les volatilités pour déterminer les seuils
    sorted_vols = sorted(volatilites)
    n = len(sorted_vols)
    
    # Seuils basés sur les percentiles 33 et 66
    threshold_low = sorted_vols[int(n * 0.33)]
    threshold_high = sorted_vols[int(n * 0.66)]
    
    # Classification des régimes
    regimes = []
    for vol in volatilites:
        if vol <= threshold_low:
            regimes.append(0)  # Faible
        elif vol <= threshold_high:
            regimes.append(1)  # Intermédiaire
        else:
            regimes.append(2)  # Élevé
    
    # Statistiques par régime
    regime_stats = {
        0: {'nom': 'Faible', 'count': 0, 'vols': [], 'color': '#28a745'},
        1: {'nom': 'Intermédiaire', 'count': 0, 'vols': [], 'color': '#ffc107'},
        2: {'nom': 'Élevé', 'count': 0, 'vols': [], 'color': '#dc3545'}
    }
    
    for i, regime in enumerate(regimes):
        regime_stats[regime]['count'] += 1
        regime_stats[regime]['vols'].append(volatilites[i])
    
    # Calculer moyennes et écarts-types par régime
    for regime in regime_stats:
        if regime_stats[regime]['vols']:
            vols = regime_stats[regime]['vols']
            regime_stats[regime]['mean'] = round(sum(vols) / len(vols), 2)
            regime_stats[regime]['min'] = round(min(vols), 2)
            regime_stats[regime]['max'] = round(max(vols), 2)
        else:
            regime_stats[regime]['mean'] = 0
            regime_stats[regime]['min'] = 0
            regime_stats[regime]['max'] = 0
        del regime_stats[regime]['vols']  # Supprimer la liste pour alléger la réponse
    
    # Matrice de transition (probabilités de passer d'un régime à un autre)
    transition_matrix = [[0, 0, 0], [0, 0, 0], [0, 0, 0]]
    transition_counts = [[0, 0, 0], [0, 0, 0], [0, 0, 0]]
    
    for i in range(1, len(regimes)):
        prev_regime = regimes[i - 1]
        curr_regime = regimes[i]
        transition_counts[prev_regime][curr_regime] += 1
    
    # Convertir en probabilités
    for i in range(3):
        total = sum(transition_counts[i])
        if total > 0:
            for j in range(3):
                transition_matrix[i][j] = round(transition_counts[i][j] / total * 100, 1)
    
    # Régime actuel
    current_regime = regimes[-1] if regimes else 0
    current_volatility = round(volatilites[-1], 2) if volatilites else 0
    
    # Données pour le graphique
    chart_data = []
    for i in range(len(volatilites)):
        chart_data.append({
            'date': vol_dates[i],
            'volatility': round(volatilites[i], 2),
            'regime': regimes[i]
        })
    
    return JsonResponse({
        'fcp_name': fcp_name,
        'window': window,
        'chart_data': chart_data,
        'regime_stats': regime_stats,
        'transition_matrix': transition_matrix,
        'current_regime': current_regime,
        'current_volatility': current_volatility,
        'thresholds': {
            'low': round(threshold_low, 2),
            'high': round(threshold_high, 2)
        },
        'total_observations': len(volatilites)
    })


def api_rolling_metrics(request):
    """API pour les métriques glissantes (Sharpe et Beta)"""
    fcp_name = request.GET.get('fcp')
    window = int(request.GET.get('window', 20))
    benchmark = request.GET.get('benchmark', 'FCP ACTIONS PERFORMANCES')
    
    # Limiter la fenêtre entre 10 et 60
    window = max(10, min(60, window))
    
    if not fcp_name:
        return JsonResponse({'error': 'FCP non spécifié'}, status=400)
    
    vl_model = get_vl_model(fcp_name)
    if not vl_model:
        return JsonResponse({'error': 'FCP non trouvé'}, status=404)
    
    vl_queryset = vl_model.objects.all().order_by('date')
    
    if not vl_queryset.exists() or vl_queryset.count() < window + 10:
        return JsonResponse({'error': 'Données insuffisantes'}, status=400)
    
    # Calculer les rendements quotidiens du FCP
    vl_list = list(vl_queryset.values('date', 'valeur'))
    rendements = []
    dates = []
    
    for i in range(1, len(vl_list)):
        ret = (float(vl_list[i]['valeur']) / float(vl_list[i-1]['valeur']) - 1) * 100
        rendements.append(ret)
        dates.append(vl_list[i]['date'])
    
    # Calculer les rendements du benchmark
    benchmark_returns = {}
    benchmark_model = get_vl_model(benchmark)
    if benchmark_model and benchmark != fcp_name:
        bench_vl = list(benchmark_model.objects.all().order_by('date').values('date', 'valeur'))
        for i in range(1, len(bench_vl)):
            ret = (float(bench_vl[i]['valeur']) / float(bench_vl[i-1]['valeur']) - 1) * 100
            benchmark_returns[bench_vl[i]['date']] = ret
    
    # Taux sans risque journalier (utiliser la constante)
    # Note: rf_daily est déjà défini en constante globale RISK_FREE_RATE_DAILY
    
    # Calculer les métriques glissantes
    rolling_sharpe = []
    rolling_beta = []
    rolling_dates = []
    
    for i in range(window - 1, len(rendements)):
        window_returns = rendements[i - window + 1:i + 1]
        window_dates_list = dates[i - window + 1:i + 1]
        
        # Sharpe Ratio glissant
        mean_ret = sum(window_returns) / len(window_returns)
        variance = sum((r - mean_ret) ** 2 for r in window_returns) / (len(window_returns) - 1) if len(window_returns) > 1 else sum((r - mean_ret) ** 2 for r in window_returns)
        std_ret = variance ** 0.5
        
        if std_ret > 0:
            sharpe = ((mean_ret - RISK_FREE_RATE_DAILY) / std_ret) * ANNUALIZATION_FACTOR
        else:
            sharpe = 0
        
        rolling_sharpe.append(round(sharpe, 3))
        rolling_dates.append(dates[i].strftime('%Y-%m-%d'))
        
        # Beta glissant (si benchmark disponible)
        if benchmark_returns:
            bench_window = [benchmark_returns.get(d, 0) for d in window_dates_list]
            
            # Vérifier qu'on a assez de données benchmark
            if sum(1 for b in bench_window if b != 0) > window * 0.5:
                mean_bench = sum(bench_window) / len(bench_window)
                
                # Covariance et variance du benchmark
                covariance = sum((window_returns[k] - mean_ret) * (bench_window[k] - mean_bench) 
                                for k in range(len(window_returns))) / len(window_returns)
                var_bench = sum((b - mean_bench) ** 2 for b in bench_window) / len(bench_window)
                
                if var_bench > 0:
                    beta = covariance / var_bench
                else:
                    beta = 1
                
                rolling_beta.append(round(beta, 3))
            else:
                rolling_beta.append(None)
        else:
            rolling_beta.append(None)
    
    return JsonResponse({
        'fcp_name': fcp_name,
        'window': window,
        'benchmark': benchmark if benchmark_returns else None,
        'dates': rolling_dates,
        'sharpe': rolling_sharpe,
        'beta': rolling_beta,
        'current_sharpe': rolling_sharpe[-1] if rolling_sharpe else None,
        'current_beta': rolling_beta[-1] if rolling_beta and rolling_beta[-1] is not None else None
    })


def api_tail_risk(request):
    """API pour l'analyse du Tail Risk (événements extrêmes)"""
    fcp_name = request.GET.get('fcp')
    
    if not fcp_name:
        return JsonResponse({'error': 'FCP non spécifié'}, status=400)
    
    vl_model = get_vl_model(fcp_name)
    if not vl_model:
        return JsonResponse({'error': 'FCP non trouvé'}, status=404)
    
    vl_queryset = vl_model.objects.all().order_by('date')
    
    if not vl_queryset.exists() or vl_queryset.count() < 30:
        return JsonResponse({'error': 'Données insuffisantes'}, status=400)
    
    # Calculer les rendements quotidiens
    vl_list = list(vl_queryset.values('date', 'valeur'))
    rendements = []
    dates = []
    
    for i in range(1, len(vl_list)):
        ret = (float(vl_list[i]['valeur']) / float(vl_list[i-1]['valeur']) - 1) * 100
        rendements.append(ret)
        dates.append(vl_list[i]['date'].strftime('%Y-%m-%d'))
    
    # Calculer moyenne et écart-type
    n = len(rendements)
    mean_ret = sum(rendements) / n
    variance = sum((r - mean_ret) ** 2 for r in rendements) / n
    std_ret = variance ** 0.5
    
    # Seuils sigma
    sigma_1 = mean_ret - std_ret
    sigma_2 = mean_ret - 2 * std_ret
    sigma_3 = mean_ret - 3 * std_ret
    
    # Compteurs d'événements extrêmes (pertes)
    losses_1_sigma = []  # Pertes > 1σ
    losses_2_sigma = []  # Pertes > 2σ
    losses_3_sigma = []  # Pertes > 3σ
    
    for i, ret in enumerate(rendements):
        if ret < sigma_1:
            losses_1_sigma.append({'date': dates[i], 'return': round(ret, 3)})
            if ret < sigma_2:
                losses_2_sigma.append({'date': dates[i], 'return': round(ret, 3)})
                if ret < sigma_3:
                    losses_3_sigma.append({'date': dates[i], 'return': round(ret, 3)})
    
    # Statistiques
    total_days = len(rendements)
    
    # Distribution théorique vs réelle (règle empirique: 68-95-99.7)
    theoretical_1_sigma = 15.87  # % attendu au-delà de 1σ (en dessous de la moyenne)
    theoretical_2_sigma = 2.28   # % attendu au-delà de 2σ
    theoretical_3_sigma = 0.13   # % attendu au-delà de 3σ
    
    actual_1_sigma = (len(losses_1_sigma) / total_days) * 100
    actual_2_sigma = (len(losses_2_sigma) / total_days) * 100
    actual_3_sigma = (len(losses_3_sigma) / total_days) * 100
    
    # Top 10 pires jours
    worst_days = sorted(zip(dates, rendements), key=lambda x: x[1])[:10]
    worst_days_list = [{'date': d, 'return': round(r, 3)} for d, r in worst_days]
    
    # Top 10 meilleurs jours
    best_days = sorted(zip(dates, rendements), key=lambda x: x[1], reverse=True)[:10]
    best_days_list = [{'date': d, 'return': round(r, 3)} for d, r in best_days]
    
    # Calculer le rapport gain/perte des événements extrêmes
    extreme_losses = [r for r in rendements if r < sigma_2]
    extreme_gains = [r for r in rendements if r > mean_ret + 2 * std_ret]
    
    avg_extreme_loss = sum(extreme_losses) / len(extreme_losses) if extreme_losses else 0
    avg_extreme_gain = sum(extreme_gains) / len(extreme_gains) if extreme_gains else 0
    
    return JsonResponse({
        'fcp_name': fcp_name,
        'total_days': total_days,
        'statistics': {
            'mean': round(mean_ret, 4),
            'std': round(std_ret, 4),
            'sigma_1_threshold': round(sigma_1, 4),
            'sigma_2_threshold': round(sigma_2, 4),
            'sigma_3_threshold': round(sigma_3, 4)
        },
        'tail_analysis': {
            '1_sigma': {
                'count': len(losses_1_sigma),
                'frequency': round(actual_1_sigma, 2),
                'expected': theoretical_1_sigma,
                'ratio': round(actual_1_sigma / theoretical_1_sigma, 2) if theoretical_1_sigma > 0 else 0,
                'events': losses_1_sigma[-20:]  # 20 derniers événements
            },
            '2_sigma': {
                'count': len(losses_2_sigma),
                'frequency': round(actual_2_sigma, 2),
                'expected': theoretical_2_sigma,
                'ratio': round(actual_2_sigma / theoretical_2_sigma, 2) if theoretical_2_sigma > 0 else 0,
                'events': losses_2_sigma[-10:]  # 10 derniers événements
            },
            '3_sigma': {
                'count': len(losses_3_sigma),
                'frequency': round(actual_3_sigma, 2),
                'expected': theoretical_3_sigma,
                'ratio': round(actual_3_sigma / theoretical_3_sigma, 2) if theoretical_3_sigma > 0 else 0,
                'events': losses_3_sigma  # Tous les événements (rares)
            }
        },
        'worst_days': worst_days_list,
        'best_days': best_days_list,
        'extreme_stats': {
            'avg_extreme_loss': round(avg_extreme_loss, 3),
            'avg_extreme_gain': round(avg_extreme_gain, 3),
            'extreme_losses_count': len(extreme_losses),
            'extreme_gains_count': len(extreme_gains)
        }
    })


def api_calendar_data(request):
    """API pour les données du calendrier de performance"""
    fcp_name = request.GET.get('fcp')
    
    if not fcp_name:
        return JsonResponse({'error': 'FCP non spécifié'}, status=400)
    
    vl_model = get_vl_model(fcp_name)
    if not vl_model:
        return JsonResponse({'error': 'FCP non trouvé'}, status=404)
    
    vl_queryset = vl_model.objects.all().order_by('date')
    
    if not vl_queryset.exists() or vl_queryset.count() < 30:
        return JsonResponse({'error': 'Données insuffisantes'}, status=400)
    
    # Calculer les rendements quotidiens
    vl_list = list(vl_queryset.values('date', 'valeur'))
    daily_returns = {}
    
    for i in range(1, len(vl_list)):
        date = vl_list[i]['date']
        ret = (float(vl_list[i]['valeur']) / float(vl_list[i-1]['valeur']) - 1) * 100
        daily_returns[date] = ret
    
    # 1. Heatmap mensuelle (performance par mois)
    monthly_data = {}
    monthly_vl_start = {}
    
    for vl in vl_list:
        date = vl['date']
        year = date.year
        month = date.month
        key = f"{year}-{month:02d}"
        
        if key not in monthly_vl_start:
            monthly_vl_start[key] = float(vl['valeur'])
        monthly_data[key] = float(vl['valeur'])
    
    # Calculer les rendements mensuels
    monthly_returns = {}
    prev_month_end = None
    for key in sorted(monthly_data.keys()):
        if prev_month_end is not None:
            monthly_returns[key] = ((monthly_data[key] / prev_month_end) - 1) * 100
        prev_month_end = monthly_data[key]
    
    # Format heatmap
    heatmap_data = []
    for key, ret in monthly_returns.items():
        year, month = key.split('-')
        heatmap_data.append({
            'year': int(year),
            'month': int(month),
            'return': round(ret, 2)
        })
    
    # 2. Performance par jour de la semaine
    weekday_stats = {i: [] for i in range(5)}  # 0=Lundi, 4=Vendredi
    weekday_names = ['Lundi', 'Mardi', 'Mercredi', 'Jeudi', 'Vendredi']
    
    for date, ret in daily_returns.items():
        weekday = date.weekday()
        if weekday < 5:  # Exclure weekends
            weekday_stats[weekday].append(ret)
    
    weekday_analysis = []
    for i in range(5):
        if weekday_stats[i]:
            returns = weekday_stats[i]
            mean = sum(returns) / len(returns)
            positive = sum(1 for r in returns if r > 0)
            negative = sum(1 for r in returns if r < 0)
            weekday_analysis.append({
                'day': weekday_names[i],
                'dayIndex': i,
                'mean': round(mean, 4),
                'count': len(returns),
                'positive': positive,
                'negative': negative,
                'win_rate': round(positive / len(returns) * 100, 1),
                'best': round(max(returns), 2),
                'worst': round(min(returns), 2)
            })
    
    # 3. Saisonnalité - Performance par mois (historique tous les mois)
    month_stats = {i: [] for i in range(1, 13)}
    month_names = ['Janvier', 'Février', 'Mars', 'Avril', 'Mai', 'Juin',
                   'Juillet', 'Août', 'Septembre', 'Octobre', 'Novembre', 'Décembre']
    
    for key, ret in monthly_returns.items():
        month = int(key.split('-')[1])
        month_stats[month].append(ret)
    
    seasonality = []
    for month in range(1, 13):
        if month_stats[month]:
            returns = month_stats[month]
            mean = sum(returns) / len(returns)
            positive = sum(1 for r in returns if r > 0)
            seasonality.append({
                'month': month_names[month - 1],
                'monthIndex': month,
                'mean': round(mean, 2),
                'count': len(returns),
                'positive': positive,
                'negative': len(returns) - positive,
                'win_rate': round(positive / len(returns) * 100, 1),
                'best': round(max(returns), 2),
                'worst': round(min(returns), 2)
            })
    
    # Trier pour trouver les meilleurs/pires mois
    sorted_seasonality = sorted(seasonality, key=lambda x: x['mean'], reverse=True)
    best_months = sorted_seasonality[:3]
    worst_months = sorted_seasonality[-3:]
    
    # 4. Données pour heatmap quotidienne (dernières 52 semaines)
    from datetime import timedelta
    if vl_list:
        last_date = vl_list[-1]['date']
        start_date = last_date - timedelta(days=365)
        
        daily_heatmap = []
        for date, ret in daily_returns.items():
            if date >= start_date:
                daily_heatmap.append({
                    'date': date.strftime('%Y-%m-%d'),
                    'weekday': date.weekday(),
                    'week': date.isocalendar()[1],
                    'year': date.year,
                    'return': round(ret, 3)
                })
    
    return JsonResponse({
        'fcp_name': fcp_name,
        'monthly_heatmap': heatmap_data,
        'weekday_analysis': weekday_analysis,
        'seasonality': seasonality,
        'best_months': best_months,
        'worst_months': worst_months,
        'daily_heatmap': daily_heatmap
    })
